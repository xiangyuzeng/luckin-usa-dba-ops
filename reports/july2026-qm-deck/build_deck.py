#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
US Quality Management — Monthly Report, July 2026 (PPTX).
Mirrors the June 2026 deck structure (5 sections) using real July data:
  01/02 Supplier & Warehouse  <- luckyus_scm_srm.t_pqnc (MCP-SSE pull, raw/*.json)
  03    Store Audit           <- /app/reports/july2026-qa-inspection (CSV + datapack)
  04/05 Complaint & EHS       <- no system of record available; explicit placeholders

Design tokens are lifted from the June 2026 deck so the two months read as one
series. Extracted from `reports/QM Monthly Report- 2026 Jun.pptx`:
  font    Microsoft YaHei everywhere (June: 496 runs)
  title   20pt bold 1E2D41   subtitle/note 11pt 5F6978   page no. 8pt 5F6978
  KPI     card F5F7FA / border E1E1E1 / value 24pt bold in a status colour /
          label 9.8pt bold 1E2D41 / note 8.3pt 5F6978  (no left accent bar)
  status  info 005EB8  good 288C5A  warn EB9119  crit C82D2D  serious C0531C
  chart   primary 1F3864, secondary 9DC3E6, then 288C5A / EB9119
  rules   body border C9C9C9, card border E1E1E1, gridline E8E8E8
Note: June itself mixes two looks — pp.1-12/31-36 hand-built, pp.13-30 (QA) already
carry the title rule + footer used here. We keep the rule/footer on every page.
"""
import csv, json, re
from collections import defaultdict, Counter
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn

HERE = Path(__file__).resolve().parent
RAW  = HERE / "raw"
QA   = Path("/app/reports/july2026-qa-inspection")
OUT  = HERE / "output"; OUT.mkdir(exist_ok=True)
OUT_PATH = OUT / "QM Monthly Report- 2026 Jul.pptx"

# ---------------------------------------------------------------- palette (June tokens)
NAVY   = "1F3864"   # brand navy — fills: bands, table headers, dividers, footer
TITLE  = "1E2D41"   # title / KPI-label ink
INK    = "222222"
INK2   = "5F6978"   # subtitle, notes, secondary text
MUTED  = "5F6978"
AXIS   = "888888"   # chart tick labels
WHITE  = "FFFFFF"
RULE   = "C9C9C9"   # body-block border
CARD_L = "E1E1E1"   # KPI-card border
GRID   = "E8E8E8"
FILL_L = "F5F7FA"   # KPI-card fill
SURF   = "FFFFFF"   # body-block fill
FONT   = "Microsoft YaHei"
# chart slots
C1, C2, C3, C4 = "1F3864", "9DC3E6", "288C5A", "EB9119"
# sequential navy ramp (ordinal: severity / grade)
SEQ = ["1F3864", "1F4E96", "2E75B6", "9DC3E6", "D6E4F5"]
# status (KPI values / accents, always alongside a text label)
INFO = "005EB8"
GOOD, WARN, SERIOUS, CRIT = "288C5A", "EB9119", "C0531C", "C82D2D"

def rgb(h): return RGBColor.from_string(h)


def set_font(run, name=FONT):
    """Set the latin AND east-Asian typeface on a run.

    The east-Asian face must be a child element <a:ea typeface="..."/> placed
    directly after <a:latin>. Writing it as an `eastAsianTypeface` ATTRIBUTE on
    <a:rPr> produces well-formed but schema-invalid XML: PowerPoint refuses the
    file with "found a problem with content" while python-pptx and most viewers
    accept it silently. Verified against the June deck, which uses <a:ea>.
    """
    run.font.name = name                       # writes <a:latin typeface="..."/>
    rPr = run.font._rPr
    latin = rPr.find(qn("a:latin"))
    if latin is None:                          # nothing to anchor to; latin alone suffices
        return
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        latin.addnext(ea)                      # schema order: ... latin, ea, cs ...
    ea.set("typeface", name)

# ---------------------------------------------------------------- data
def L(p): return json.loads(Path(p).read_text(encoding="utf-8"))
PJ   = L(RAW/"pqnc_july.json")
PJUN = L(RAW/"pqnc_june.json")
PTR  = L(RAW/"pqnc_trend.json")
SUPP = {s["supplier_mid"]: s["supplier_name"] for s in L(RAW/"suppliers.json")}
# store / goods names for July rows — emitted by verify_vs_xlsx.py from the official
# SRM export `reports/PQNC 2026-07.xlsx` after it is proven row-identical to the pull.
ENR  = L(RAW/"xlsx_enrich.json")
def store(r):  return ENR.get(r["pqnc_no"], {}).get("store", "—")
def goods(r):  return ENR.get(r["pqnc_no"], {}).get("goods", "—")

PACK = L(QA/"july2026_qa_datapack.json")
# Orkin per-store pest service reports -> ../july2026-pest-service/build_pest_pack.py
PEST = L(Path("/app/reports/july2026-pest-service/july2026_pest_datapack.json"))
PD   = PEST["derived"]
DER  = L(QA/"derived.json")
def load_csv(p):
    with open(p, encoding="utf-8-sig") as f: return list(csv.DictReader(f))
SUM  = load_csv(QA/"july2026_inspection_summary.csv")
ITEM = load_csv(QA/"july2026_inspection_items.csv")
for r in SUM:
    for k in ("inspection_id","adjusted_total_score","original_total_score",
              "S_count","M_count","G_count","L_count","original_total_deduction"):
        r[k] = int(r[k]) if r[k] not in ("","None",None) else 0
for r in ITEM:
    r["deduction"] = int(r["deduction"]); r["inspection_id"] = int(r["inspection_id"])

NAME = {r["store_code"]: r["store_name"] for r in SUM}
def sn(c): return NAME.get(c, c)

PRIMARY = DER["primary_by_store"]; JUNP = DER["june_primary"]
NPRIM = len(PRIMARY)
PRIM_IIDS = set(DER["primary_iids"])
PITEMS = [it for it in ITEM if it["inspection_id"] in PRIM_IIDS]

# ---- QA-scope (deck convention: latest QA audit per store) ----
qa_rows = [r for r in SUM if r["inspection_type"] == "QA审计"]
_by = defaultdict(list)
for r in qa_rows: _by[r["store_code"]].append(r)
QA_LATEST = {c: sorted(v, key=lambda r: (r["inspection_date"], r["inspection_id"]))[-1]
             for c, v in _by.items()}
QA_N, QA_STORES = len(qa_rows), len(QA_LATEST)
qa_adj = [r["adjusted_total_score"] for r in QA_LATEST.values()]
qa_org = [r["original_total_score"] for r in QA_LATEST.values()]
QA_AVG, QA_AVG_ORG = round(sum(qa_adj)/len(qa_adj),1), round(sum(qa_org)/len(qa_org),1)
QA_S = sum(r["S_count"] for r in qa_rows); QA_M = sum(r["M_count"] for r in qa_rows)

def grade(x): return "A+" if x>=94 else ("A" if x>=87 else ("B" if x>=80 else "C"))
GA = Counter(grade(x) for x in qa_adj); GO = Counter(grade(x) for x in qa_org)

FM = PACK["cover"]["findings_full_month"]; PM = PACK["cover"]["findings_primary"]
TC = PACK["cover"]["type_counts"]; AP = PACK["cover"]["appeals"]
REPEAT = PACK["sX_repeat_offenders"]
MA = PACK["module_agg_primary"]
# stores whose PRIMARY inspection carried a Sinks-and-Pipes S finding
PRIM_PIPE = sorted({x["store"] for x in PACK["s3_2"] if x["sub_item"] == "Sinks and Pipes"})

# ---- PQNC derived ----
RESP = {1:"Supplier", 2:"Warehouse", 3:"Store", 4:"Joint", 5:"Unknown/reject", 6:"Unknown/reject"}
def resp(r): return RESP.get(r["resp_code"], "Unknown/reject")
def val(r):  return float(r["value_amount"] or 0)
PJ_RESP = Counter(resp(r) for r in PJ); PJUN_RESP = Counter(resp(r) for r in PJUN)
TYPE = {"0003":"Food Safety Issue", "0004":"General Defect", "0001":"Sensory Abnormal", "0002":"Other Unclear"}
def ptype(r): return TYPE.get(r["one_pqnc_type_code"], "Unclassified")
PJ_TYPE = Counter(ptype(r) for r in PJ); PJUN_TYPE = Counter(ptype(r) for r in PJUN)
PJ_VAL = round(sum(val(r) for r in PJ), 2)
N_JUL, N_JUN = len(PJ), len(PJUN)
FS_JUL = PJ_TYPE.get("Food Safety Issue", 0); FS_JUN = PJUN_TYPE.get("Food Safety Issue", 0)
WH_JUL = PJ_RESP.get("Warehouse",0) + PJ_RESP.get("Joint",0)
WH_JUN = PJUN_RESP.get("Warehouse",0) + PJUN_RESP.get("Joint",0)
SP_JUL, SP_JUN = PJ_RESP.get("Supplier",0), PJUN_RESP.get("Supplier",0)
# July's joint-responsibility bucket is one single event: the Cream-O-Land fat-free
# milk spoilage cluster, judged 2026-08-03 (after the month closed).
JOINT    = [r for r in PJ if resp(r) == "Joint"]
JOINT_VAL = sum(val(r) for r in JOINT)

def wh_issue(r):
    """Warehouse/joint issue class — same keyword rule applied to both months."""
    t = (r["problem_description"] or "").lower()
    if any(k in t for k in ["expire","expiration","short shelf","out of date","best by"]): return "过期/短保"
    if any(k in t for k in ["temp","frozen","thaw","defrost","warm","room temperature","cold"]): return "失温/错误温区"
    if any(k in t for k in ["leak","damag","dent","broken","rip","crack","busted","seal","hole","spill"]): return "破损变形/泄漏"
    if any(k in t for k in ["spoil","sour","mold","bad smell","rancid"]): return "变质"
    return "其他"

def goods_cat(r):
    t = ((r["problem_description"] or "") + " " + (r["factory_name"] or "")).lower()
    if any(k in t for k in ["milk","cream","dairy","o land","o-land","coconut"]): return "奶类 Milk"
    if any(k in t for k in ["coffee","bean","roast","espresso"]):                 return "咖啡豆 Coffee"
    if any(k in t for k in ["croissant","cookie","pastry","cheese","sandwich"]):  return "轻食 Light food"
    if any(k in t for k in ["cup","bottle","package","bag","label","can","lid","seal","shelf"]): return "包材/器具 Packaging"
    return "其他 Other"

# ---------------------------------------------------------------- deck helpers
prs = Presentation()
prs.slide_width  = Emu(12192000)   # 13.333in
prs.slide_height = Emu(6858000)    # 7.5in
BLANK = prs.slide_layouts[6]

def S(): return prs.slides.add_slide(BLANK)

def tb(sl, x, y, w, h, text, size=11, bold=False, color=INK, align=PP_ALIGN.LEFT,
       anchor=MSO_ANCHOR.TOP, spacing=1.0, wrap=True):
    box = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Emu(0); tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = rgb(color)
        set_font(r)
    return box

def rect(sl, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, lw=0.75):
    s = sl.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = rgb(fill)
    else:    s.fill.background()
    if line: s.line.color.rgb = rgb(line); s.line.width = Pt(lw)
    else:    s.line.fill.background()
    s.shadow.inherit = False
    return s

def title_bar(sl, title, sub=None, page=None):
    tb(sl, 0.55, 0.26, 12.2, 0.5, title, size=20, bold=True, color=TITLE)
    rect(sl, 0.55, 0.86, 12.2, 0.022, fill=NAVY)
    if sub: tb(sl, 0.55, 0.94, 12.2, 0.3, sub, size=11, color=INK2)
    tb(sl, 9.3, 7.02, 3.5, 0.25, "瑞幸咖啡 · 北美质量保障部", size=10, color=NAVY, align=PP_ALIGN.RIGHT)
    if page is not None:
        tb(sl, 12.5, 7.03, 0.4, 0.25, str(page), size=8, color=MUTED, align=PP_ALIGN.RIGHT)

def kpi(sl, x, y, w, h, value, label, note=None, accent=INFO, vsize=24):
    rect(sl, x, y, w, h, fill=FILL_L, line=CARD_L)
    tb(sl, x+0.16, y+0.12, w-0.28, 0.45, value, size=vsize, bold=True, color=accent)
    tb(sl, x+0.16, y+0.12+0.44, w-0.28, 0.42, label, size=9.8, bold=True, color=TITLE, spacing=0.95)
    if note: tb(sl, x+0.16, y+h-0.28, w-0.28, 0.24, note, size=8.3, color=MUTED)

def band(sl, x, y, w, title):
    rect(sl, x, y, w, 0.34, fill=NAVY)
    tb(sl, x+0.12, y+0.06, w-0.24, 0.24, title, size=10.5, bold=True, color=WHITE)

def body(sl, x, y, w, h, lines, size=10, fill=SURF, line=RULE, spacing=1.25):
    rect(sl, x, y, w, h, fill=fill, line=line)
    tb(sl, x+0.16, y+0.13, w-0.32, h-0.26, lines, size=size, color=INK, spacing=spacing)

def table(sl, x, y, w, h, headers, rows, widths=None, fsize=9, hsize=9,
          align_center=None, row_colors=None):
    t = sl.shapes.add_table(len(rows)+1, len(headers), Inches(x), Inches(y),
                            Inches(w), Inches(h)).table
    if widths:
        tot = sum(widths)
        for i, cw in enumerate(widths): t.columns[i].width = Emu(int(Inches(w)*cw/tot))
    for j, hd in enumerate(headers):
        c = t.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = rgb(NAVY)
        c.margin_left = c.margin_right = Emu(45720); c.margin_top = c.margin_bottom = Emu(18288)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(hd); r.font.size = Pt(hsize); r.font.bold = True
        r.font.color.rgb = rgb(WHITE); set_font(r)
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            c = t.cell(i+1, j)
            c.fill.solid()
            c.fill.fore_color.rgb = rgb((row_colors or {}).get(i) or (FILL_L if i % 2 else WHITE))
            c.margin_left = c.margin_right = Emu(45720); c.margin_top = c.margin_bottom = Emu(13716)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if (align_center and j in align_center) else PP_ALIGN.LEFT
            txt, opt = (v if isinstance(v, tuple) else (v, {}))
            r = p.add_run(); r.text = str(txt)
            r.font.size = Pt(opt.get("size", fsize)); r.font.bold = opt.get("bold", False)
            r.font.color.rgb = rgb(opt.get("color", INK)); set_font(r)
    return t

def style_chart(ch, colors, labels=True, legend=True, num_fmt="0", lsize=8, gap=60):
    ch.font.size = Pt(9); ch.font.name = FONT; ch.font.color.rgb = rgb(INK2)
    if legend and len(ch.plots[0].series) > 1:
        ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False; ch.legend.font.size = Pt(8.5)
    else:
        ch.has_legend = False
    try: ch.plots[0].gap_width = gap
    except Exception: pass
    for i, ser in enumerate(ch.plots[0].series):
        col = colors[i % len(colors)]
        try:
            ser.format.fill.solid(); ser.format.fill.fore_color.rgb = rgb(col)
            ser.format.line.color.rgb = rgb(col); ser.format.line.width = Pt(2)
        except Exception: pass
    if labels:
        pl = ch.plots[0]; pl.has_data_labels = True
        dl = pl.data_labels; dl.font.size = Pt(lsize); dl.font.color.rgb = rgb(INK)
        dl.number_format = num_fmt; dl.number_format_is_linked = False
        try: dl.position = XL_LABEL_POSITION.OUTSIDE_END
        except Exception:
            try: dl.position = XL_LABEL_POSITION.CENTER
            except Exception: pass
    try:
        ch.value_axis.has_major_gridlines = True
        gl = ch.value_axis.major_gridlines.format.line
        gl.color.rgb = rgb(GRID); gl.width = Pt(0.75)
        ch.value_axis.format.line.color.rgb = rgb(RULE)
        ch.category_axis.format.line.color.rgb = rgb(RULE)
        ch.value_axis.tick_labels.font.size = Pt(8)
        ch.category_axis.tick_labels.font.size = Pt(8)
    except Exception: pass
    return ch

def chart(sl, kind, x, y, w, h, cats, series, colors, **kw):
    cd = CategoryChartData(); cd.categories = cats
    for nm, vals in series: cd.add_series(nm, vals)
    gf = sl.shapes.add_chart(kind, Inches(x), Inches(y), Inches(w), Inches(h), cd)
    return style_chart(gf.chart, colors, **kw)

def placeholder(sl, x, y, w, h, title, need, ref):
    rect(sl, x, y, w, h, fill="FFF8E8", line=WARN, lw=1.25)
    tb(sl, x+0.25, y+0.22, w-0.5, 0.3, f"⚠ {title}", size=13, bold=True, color="9A6B00")
    tb(sl, x+0.25, y+0.62, w-0.5, h-0.9,
       need + "\n\n" + ref, size=10.5, color=INK, spacing=1.35)

ASSETS = HERE / "assets"     # cover / divider / thank-you artwork lifted from the June deck
def pic(sl, name, x, y, w, h):
    return sl.shapes.add_picture(str(ASSETS/f"{name}.png"), Inches(x), Inches(y), Inches(w), Inches(h))

def section_head(sl, num, title, x, y, w, h, size=32, color=None):
    """June-style section opener: number label above a large bilingual title."""
    if num: tb(sl, x, y-0.52, w, 0.34, num, size=14, bold=True, color=NAVY)
    tb(sl, x, y, w, h, title, size=size, bold=True, color=color or TITLE, spacing=1.18)

ATTACH_NOTE = ("※ 本页对应的现场实拍附件存于对象存储（工单 attachment_url 仅存路径），"
               "未纳入本次数据拉取；如需配图请从 SRM / 稽核系统导出后插入。")
def attach_note(sl, x=0.55, w=12.2):
    """Drop the note just under the lowest block on the slide (June had photos here).
    If the slide is already full, fall back to the footer row beside the org name."""
    bot = max((sh.top + sh.height) / 914400 for sh in sl.shapes
              if (sh.top + sh.height) / 914400 < 6.95)
    if bot + 0.06 <= 6.74:
        tb(sl, x, bot + 0.06, w, 0.22, ATTACH_NOTE, size=8, color=MUTED)
    else:
        tb(sl, x, 7.03, 8.6, 0.22, ATTACH_NOTE, size=7.5, color=MUTED)

PG = [0]
def page(): PG[0] += 1; return PG[0]

# ================================================================= 1 COVER
s = S()
pic(s, "cover", 0, -0.02, 4.99, 7.49)
tb(s, 6.57, 2.70, 5.53, 1.9, "US Quality Management\n北美质量管理部", size=32, bold=True,
   color=TITLE, spacing=1.2)
rect(s, 6.60, 4.78, 1.4, 0.035, fill=NAVY)
tb(s, 6.57, 5.08, 5.53, 0.6, "Jul 2026", size=32, bold=True, color=NAVY)
tb(s, 6.60, 5.86, 5.53, 0.35, "2026 年 7 月度质量报告", size=15, color=INK2)
tb(s, 6.60, 6.30, 5.53, 0.3, "编制：曾翔宇　|　日期：2026-08-04", size=10.5, color=MUTED)

# ================================================================= 2 AGENDA
s = S(); title_bar(s, "7 月整体质量表现", "Monthly Quality Performance Overview · Jul 2026", page())
AG = [("01", "Supplier & Materials Performance 供应链", 0.9, 1.35, WARN,
       f"PQNC {N_JUL} 起 · 供应商责任 {SP_JUL} 起"),
      ("02", "Warehouse Quality Performance 仓配", 6.85, 1.35, WARN,
       f"仓储·共担 {WH_JUL} 起（6 月 {WH_JUN} 起，{(WH_JUL-WH_JUN)/WH_JUN*100:+.0f}%）· 含奶类变质集群 {len(JOINT)} 起"),
      ("03", "Store Audit Performance 门店稽核", 0.9, 3.0, WARN,
       "QA 稽核 19 次 / 18 家 · 均分 92.1"),
      ("04", "Customer Complaint 客户投诉", 6.85, 3.0, CRIT,
       "食安类客诉 7 起 / 6 家门店 · 4 起指向奶变质"),
      ("05", "EHS 环境健康安全", 0.9, 4.65, GOOD, "工伤 0 起 · 连续两月零工伤")]
for num, txt, x, y, col, note in AG:
    rect(s, x, y, 5.45, 1.35, fill=FILL_L)
    rect(s, x, y, 0.06, 1.35, fill=col)
    tb(s, x+0.28, y+0.18, 0.8, 0.5, num, size=22, bold=True, color=col)
    tb(s, x+1.2, y+0.22, 4.1, 0.5, txt, size=12.5, bold=True, color=NAVY, spacing=1.05)
    tb(s, x+1.2, y+0.82, 4.1, 0.35, note, size=9.5, color=INK2)
LEG = [("发现重大偏差 / 低于预期", CRIT), ("表现波动，需要密切关注", WARN), ("整体无明显异常", GOOD)]
for i, (t, c) in enumerate(LEG):
    x = 0.9 + i*4.05
    rect(s, x, 6.35, 0.16, 0.16, fill=c)
    tb(s, x+0.28, 6.33, 3.6, 0.25, t, size=9, color=INK2)

# ================================================================= 3 DIVIDER 01/02
s = S()
pic(s, "divider0102", 0, 0, 7.42, 7.50)
section_head(s, "01 / 02", "Supplier & Materials\n& Warehouse\n供应链仓配质检", 8.34, 3.08, 4.61, 2.26)

# ================================================================= 4 SUPPLIER ADMISSION (placeholder)
s = S(); title_bar(s, "01  Supplier & Materials Performance", "1. Suppliers admission / 供应商准入", page())
kpi(s, 0.55, 1.45, 2.9, 1.15, "2", "本月准入通过 / Approved", "6 月 4 家", INFO)
kpi(s, 3.65, 1.45, 2.9, 1.15, "1", "食品类 / Food", "经 QA 审批", INFO)
kpi(s, 6.75, 1.45, 2.9, 1.15, "1", "其他非食品 / Other", "经 QA 审批", INFO)
kpi(s, 9.85, 1.45, 2.9, 1.15, "0", "现场稽核 / On-site audit", "本月无准入/年度/飞行审核", MUTED)
table(s, 0.55, 2.95, 6.4, 1.6,
      ["Supplier Category", "Approved #"],
      [["Food 食品类", 1], ["Other none-food 其他非食品", 1],
       [("Total / 合计", {"bold": True}), (2, {"bold": True})]],
      widths=[5, 2], align_center=[1])
band(s, 7.15, 2.95, 5.6, "口径说明 / Scope")
body(s, 7.15, 3.29, 5.6, 1.26,
     "本页为「经 QA 审批」的准入家数。\n"
     "部分营建类供应商不经 QA 审批即可在系统建档，\n"
     "因此系统家数会高于本页。", size=9.5)
band(s, 0.55, 4.75, 12.2, "与系统主数据的差异（已核实）")
body(s, 0.55, 5.09, 12.2, 1.6,
     "系统 `t_mdm_supplier` 7 月新增 15 条供应商主数据记录，而经 QA 审批的准入仅 2 家——两者不矛盾：\n"
     "「营建 / 工程类供应商不走 QA 准入审批流程，但仍会在系统内建档」，这部分计入主数据、不计入本页。\n"
     "因此「主数据新增家数」不能直接当作「准入通过家数」使用；v1 版据此提出的数据缺口已由业务侧澄清并关闭。\n\n"
     "※ 本月仅提供通过家数；申请家数（含被否决）未单独统计，故不列「Total applied」列。", size=9.5)

# ================================================================= 5 MATERIAL ADMISSION (placeholder)
s = S(); title_bar(s, "01  Supplier & Materials Performance", "2. Materials admission / 原料准入", page())
kpi(s, 0.55, 1.5, 3.9, 1.3, "0", "本月原料准入 / Materials admitted", "6 月同为 0", INFO)
kpi(s, 4.65, 1.5, 3.9, 1.3, "0", "新增申请 / Applications", "无在途审批", INFO)
kpi(s, 8.75, 1.5, 4.0, 1.3, "连续 2 月", "无新原料引入", "6 月 0 / 7 月 0", GOOD)
table(s, 0.55, 3.1, 6.4, 2.3,
      ["Material Category", "Approved #"],
      [["Raw material 原料", 0], ["Light food 轻食", 0], ["Food contact material 食品接触", 0],
       ["Promotional items 营销物料", 0], ["Chemicals 化学品", 0], ["Other none-food 其他", 0],
       [("Total / 合计", {"bold": True}), (0, {"bold": True})]],
      widths=[5, 2], align_center=[1])
band(s, 7.15, 3.1, 5.6, "说明 / Note")
body(s, 7.15, 3.44, 5.6, 1.96,
     "7 月无新原料准入（申请 0 / 通过 0），与 6 月一致。\n\n"
     "现有原料池保持不变，本月供应商责任 PQNC 全部\n"
     "发生在既有 SKU 上，非新引入原料导致。\n\n"
     "※ 原料准入审批不在 luckyus_scm_srm 内，\n"
     "本页数据由供应链团队提供。", size=9.5)

# ================================================================= 6 PQNC SUMMARY
s = S(); title_bar(s, "Supplier & Warehouse Quality Performance",
                   "July 2026 PQNC Summary / 2026 年 7 月 PQNC 概览", page())
fs_jul, fs_jun = FS_JUL, FS_JUN
wh_jul, wh_jun = WH_JUL, WH_JUN
sp_jul, sp_jun = SP_JUL, SP_JUN
kpi(s, 0.55, 1.35, 2.9, 1.15, str(N_JUL), "PQNC reported / 月度上报",
    f"vs Jun {N_JUN}, {(N_JUL-N_JUN)/N_JUN*100:+.0f}%", GOOD)
kpi(s, 3.65, 1.35, 2.9, 1.15, str(fs_jul), "Food-safety major / 食安风险",
    f"vs Jun {fs_jun}, {(fs_jul-fs_jun)/fs_jun*100:+.0f}%", GOOD)
kpi(s, 6.75, 1.35, 2.9, 1.15, str(wh_jul), "Warehouse + Joint / 仓储·共担",
    f"vs Jun {wh_jun}, {(wh_jul-wh_jun)/wh_jun*100:+.0f}%（含集群 {len(JOINT)} 起）", WARN)
kpi(s, 9.85, 1.35, 2.9, 1.15, str(sp_jul), "Supplier resp. / 供应商责任",
    f"vs Jun {sp_jun}, {(sp_jul-sp_jun)/sp_jun*100:+.0f}%", CRIT)
MN = ["Jan","Feb","Mar","Apr","May","Jun","Jul"]
tb(s, 0.55, 2.65, 4.0, 0.22, "PQNC 月度趋势 / Monthly trend（起）", size=9, bold=True, color=INK2)
chart(s, XL_CHART_TYPE.LINE_MARKERS, 0.5, 2.9, 4.15, 2.1, MN,
      [("PQNC", [t["n"] for t in PTR])], [C1], legend=False, lsize=7.5)
tb(s, 4.85, 2.65, 3.6, 0.22, "判责结构 / Responsibility mix（起）", size=9, bold=True, color=INK2)
dn_cats = ["Supplier 供应商","Warehouse 仓储","Joint 共担","Unknown/reject 未明确"]
dn_vals = [PJ_RESP.get("Supplier",0), PJ_RESP.get("Warehouse",0),
           PJ_RESP.get("Joint",0), PJ_RESP.get("Unknown/reject",0)]
cd = CategoryChartData(); cd.categories = dn_cats; cd.add_series("Cases", dn_vals)
gf = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(4.8), Inches(2.85), Inches(3.75), Inches(2.25), cd)
ch = gf.chart; ch.font.size = Pt(8.5); ch.font.name = "Arial"
ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False; ch.legend.font.size = Pt(8)
pts = ch.plots[0].series[0].points
for i, col in enumerate([C1, C2, C3, C4]):
    pts[i].format.fill.solid(); pts[i].format.fill.fore_color.rgb = rgb(col)
    pts[i].format.line.color.rgb = rgb(WHITE); pts[i].format.line.width = Pt(2)
ch.plots[0].has_data_labels = True
ch.plots[0].data_labels.font.size = Pt(8.5); ch.plots[0].data_labels.font.bold = True
ch.plots[0].data_labels.font.color.rgb = rgb(WHITE)
tb(s, 8.75, 2.65, 4.0, 0.22, "PQNC Type / 问题类型（起）", size=9, bold=True, color=INK2)
table(s, 8.7, 2.9, 4.05, 1.5, ["PQNC type / 类型", "Cases", "占比"],
      [["Food Safety Issue / 食安", fs_jul, f"{fs_jul/N_JUL*100:.1f}%"],
       ["General Defect / 普通", PJ_TYPE.get("General Defect",0), f"{PJ_TYPE.get('General Defect',0)/N_JUL*100:.1f}%"],
       ["Unclassified / 未判定", PJ_TYPE.get("Unclassified",0), f"{PJ_TYPE.get('Unclassified',0)/N_JUL*100:.1f}%"],
       [("Total / 合计",{"bold":True}), (N_JUL,{"bold":True}), ("100%",{"bold":True})]],
      widths=[5,2,2], align_center=[1,2], fsize=8.5)
table(s, 8.7, 4.6, 4.05, 1.75, ["判责 / Responsibility", "Cases", "货值 $"],
      [["Supplier / 供应商", PJ_RESP.get("Supplier",0), f"${sum(val(r) for r in PJ if resp(r)=='Supplier'):,.2f}"],
       ["Warehouse / 仓储物流", PJ_RESP.get("Warehouse",0), f"${sum(val(r) for r in PJ if resp(r)=='Warehouse'):,.2f}"],
       ["Joint / 供应商·仓储共担", PJ_RESP.get("Joint",0), f"${JOINT_VAL:,.2f}"],
       ["Unknown / reject 未明确", PJ_RESP.get("Unknown/reject",0), f"${sum(val(r) for r in PJ if resp(r)=='Unknown/reject'):,.2f}"],
       [("Total / 合计",{"bold":True}), (N_JUL,{"bold":True}), (f"${PJ_VAL:,.2f}",{"bold":True})]],
      widths=[5,2,2], align_center=[1,2], fsize=8.5)
body(s, 0.5, 5.2, 8.05, 1.72,
     f"7 月共上报 {N_JUL} 起 PQNC，较 6 月 {N_JUN} 起微降 {abs((N_JUL-N_JUN)/N_JUN*100):.0f}%，为 3 月以来最低；货值合计 ${PJ_VAL:,.2f}。\n"
     f"食安风险仅 {fs_jul} 起（6 月 {fs_jun} 起）：7/5 Sysco 罐装产品开罐发现变色异物，已隔离留样。\n"
     f"仓储·共担 {wh_jul} 起（6 月 {wh_jun} 起，−{abs((wh_jul-wh_jun)/wh_jun*100):.0f}%）：仓储单独判责仅 {PJ_RESP.get('Warehouse',0)} 起，"
     f"其余 {len(JOINT)} 起为 7/22–7/28 Cream-O-Land 脱脂奶变质集群，8/3 判为供应商·仓储共担（详见 Case B）。\n"
     f"供应商责任 {sp_jul} 起（占 {sp_jul/N_JUL*100:.0f}%，6 月 {sp_jun} 起 {(sp_jul-sp_jun)/sp_jun*100:+.0f}%），"
     f"其中 SYSCO 23 起、集中于 Cream-O-Land 奶类效期标签与灌装量问题。\n"
     f"※ 本页数据于 2026-08-04 重算，并与 SRM 系统导出《PQNC 2026-07.xlsx》逐单核对一致（{N_JUL}/{N_JUL} 单相符）。",
     size=9)

# ================================================================= 7 SUPPLIER RESPONSIBILITY
s = S(); title_bar(s, "Supplier-Responsibility Quality Overview",
                   "供应商责任质量概况 · July 2026", page())
sup_rows = [r for r in PJ if resp(r) == "Supplier"]
sup_val = sum(val(r) for r in sup_rows)
by_sup = Counter(SUPP.get(r["supplier_mid"], "(未映射)") for r in sup_rows)
kpi(s, 0.55, 1.4, 2.6, 1.2, str(len(sup_rows)), "Supplier-resp. cases 供应商责任",
    f"vs Jun {sp_jun}, {(sp_jul-sp_jun)/sp_jun*100:+.0f}%", CRIT)
kpi(s, 3.3, 1.4, 2.6, 1.2, f"${sup_val:,.0f}", "Goods value 货值",
    f"占全月货值 {sup_val/PJ_VAL*100:.0f}%", INFO)
kpi(s, 6.05, 1.4, 2.6, 1.2, "0", "On-site supplier audits 现场稽核",
    "本月无准入/年度/飞行审核", MUTED)
tb(s, 0.55, 2.85, 4.5, 0.22, "供应商责任 PQNC 分布（起）", size=9, bold=True, color=INK2)
sup_top = by_sup.most_common(5)
chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, 0.5, 3.1, 4.4, 2.5,
      [k.split(",")[0][:16] for k,_ in sup_top], [("Cases", [v for _,v in sup_top])],
      [C2], legend=False, lsize=8)
tb(s, 5.15, 2.85, 3.4, 0.22, "供应商责任月度趋势（起）", size=9, bold=True, color=INK2)
chart(s, XL_CHART_TYPE.LINE_MARKERS, 5.1, 3.1, 3.5, 2.5, MN,
      [("Supplier", [int(t["supplier"]) for t in PTR])], [C1], legend=False, lsize=7.5)
cat_rows = defaultdict(lambda: [0, 0.0])
for r in sup_rows:
    c = goods_cat(r); cat_rows[c][0] += 1; cat_rows[c][1] += val(r)
tb(s, 8.85, 1.4, 4.0, 0.22, "按货物类别（供应商责任）※ 关键词归类", size=9, bold=True, color=INK2)
rows7 = [[k, v[0], f"${v[1]:,.2f}"] for k, v in sorted(cat_rows.items(), key=lambda kv: -kv[1][0])]
rows7.append([("Total / 合计",{"bold":True}), (len(sup_rows),{"bold":True}), (f"${sup_val:,.2f}",{"bold":True})])
table(s, 8.8, 1.68, 4.0, 1.9, ["Category / 货物类别", "Cases", "Value $"], rows7,
      widths=[5,2,2.4], align_center=[1,2], fsize=8.5)
body(s, 8.8, 3.75, 4.0, 2.9,
     f"供应商责任 PQNC {len(sup_rows)} 起，占本月总量 {len(sup_rows)/N_JUL*100:.0f}%，"
     f"较 6 月 {sp_jun} 起上升 {(sp_jul-sp_jun)/sp_jun*100:.0f}%。\n\n"
     f"SYSCO {by_sup.get('SYSCO',0)} 起：绝大多数为 Cream-O-Land 奶类——效期标签缺失/模糊 8 起、"
     f"灌装量不足 / 液位偏低 7 起、瓶体渗漏/破孔 4 起、变质酸败 3 起。\n\n"
     f"S&D Coffee 3 起：咖啡豆包装袋在箱内即已破裂散豆。\n\n"
     f"VIOBIO 3 起 / FREENOW 3 起：杯体无品牌标识、轻食成型不良。\n\n"
     f"⚠ 另有 {len(JOINT)} 起 Cream-O-Land 脱脂奶变质集群未计入本页——"
     f"该集群判为供应商·仓储共担，计入仓储页（Case B）。若并入，SYSCO 相关问题达 "
     f"{by_sup.get('SYSCO',0)+len(JOINT)} 起。\n\n"
     f"建议：对 Cream-O-Land 灌装线与效期喷码稳定性提出整改；到货抽检增加效期字迹与灌装液位项。",
     size=9)

# ================================================================= 8 WAREHOUSE
s = S(); title_bar(s, "Warehouse & Distribution Responsibility",
                   "仓储与配送责任问题 · July 2026", page())
wh_rows = [r for r in PJ if resp(r) in ("Warehouse", "Joint")]
wh_val = sum(val(r) for r in wh_rows)
wh_only = [r for r in PJ if resp(r) == "Warehouse"]
kpi(s, 0.55, 1.4, 2.7, 1.2, str(len(wh_rows)), "Warehouse + Joint 仓储·共担",
    f"vs Jun {wh_jun}, {(wh_jul-wh_jun)/wh_jun*100:+.0f}%", WARN)
kpi(s, 3.4, 1.4, 2.7, 1.2, f"${wh_val:,.0f}", "Goods value 货值",
    f"占全月货值 {wh_val/PJ_VAL*100:.0f}%", INFO)
kpi(s, 6.25, 1.4, 2.7, 1.2, str(len(JOINT)), "Joint 共担 · 奶类变质集群",
    "7/22–7/28 · 详见 Case B", CRIT)
WH_CATS = ["过期/短保","失温/错误温区","破损变形/泄漏","变质","其他"]
wh_jun_rows = [r for r in PJUN if resp(r) in ("Warehouse","Joint")]
cj = Counter(wh_issue(r) for r in wh_jun_rows); c7 = Counter(wh_issue(r) for r in wh_rows)
c7o = Counter(wh_issue(r) for r in wh_only); c7j = Counter(wh_issue(r) for r in JOINT)
tb(s, 0.55, 2.85, 6.0, 0.22, "6 月 vs 7 月 仓储责任问题分类（起）※ 关键词归类，两月同一规则",
   size=9, bold=True, color=INK2)
chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, 0.5, 3.1, 5.9, 2.6, WH_CATS,
      [("6 月 仓储·共担", [cj.get(k,0) for k in WH_CATS]),
       ("7 月 仓储", [c7o.get(k,0) for k in WH_CATS]),
       ("7 月 共担集群", [c7j.get(k,0) for k in WH_CATS])],
      [C1, C2, C3], lsize=8)
table(s, 6.6, 1.68, 6.2, 2.0, ["Issue / 问题", "判责", "Cases", "Value $"],
      [["Seal popped / 封口崩开泄漏", "仓储", 1, "$4.97"],
       ["Bottle hole / 瓶体破孔渗漏", "仓储", 1, "$4.16"],
       ["Label damaged / 标签破损缺开口拉环", "仓储", 1, "$26.18"],
       ["Missing parts / 货架配件缺件", "仓储", 1, "$398.13"],
       [("Milk spoilage cluster / 奶类变质集群",{"bold":True}), ("共担",{"bold":True}),
        (len(JOINT),{"bold":True}), (f"${JOINT_VAL:,.2f}",{"bold":True})],
       [("Total / 合计",{"bold":True}), "", (len(wh_rows),{"bold":True}), (f"${wh_val:,.2f}",{"bold":True})]],
      widths=[5.4,1.2,1.3,1.7], align_center=[1,2,3], fsize=8.5)
body(s, 6.6, 3.8, 6.2, 2.9,
     f"仓储·共担合计 {len(wh_rows)} 起（6 月 {wh_jun} 起，{(wh_jul-wh_jun)/wh_jun*100:+.0f}%），"
     f"但结构与 6 月完全不同，需拆开看：\n\n"
     f"① 仓储单独判责仅 {len(wh_only)} 起，为本年度最低——6 月「效期管控 / 冷链执行 / 配送放置」"
     f"三项专项整改见效：过期/短保 {cj.get('过期/短保',0)}→{c7o.get('过期/短保',0)} 起、"
     f"失温/错误温区 {cj.get('失温/错误温区',0)}→{c7o.get('失温/错误温区',0)} 起。"
     f"这 {len(wh_only)} 起均为包装完整性与配件缺件，其中设备仓「货架配件缺件」单笔 $398.13，"
     f"占仓储·共担货值 {398.13/wh_val*100:.0f}%，建议设备类出库增加配件清单双人核对。\n\n"
     f"② 新增 {len(JOINT)} 起「供应商·仓储共担」，全部为 Cream-O-Land 脱脂奶变质集群"
     f"（7/22–7/28，9 家门店，${JOINT_VAL:,.2f}）。判责意见：多店反馈奶变酸，供应商与仓库"
     f"双方自查均未发现异常，故共担——责任未落到单一环节，冷链链路仍是开放风险（详见 Case B）。\n\n"
     f"⚠ 另有 {len([r for r in PJ if 'defrost' in (r['problem_description'] or '').lower()])} 起"
     f"门店后厨冰箱异常除霜物料损失（${sum(val(r) for r in PJ if 'defrost' in (r['problem_description'] or '').lower()):,.2f}），"
     f"判责「未明确」，非仓储责任，需运营与设备侧介入（详见 Case C）。",
     size=8.5)

# ================================================================= 9 CASE A
s = S(); title_bar(s, "典型案例 — Milk Expiration Label / Underfill / Leakage",
                   "奶类效期标签、灌装量与渗漏 · Case A", page())
milk_all = [r for r in PJ if goods_cat(r) == "奶类 Milk"]
milk = [r for r in milk_all if resp(r) == "Supplier"]   # 共担变质集群单列 Case B
kpi(s, 0.55, 1.4, 2.7, 1.05, str(len(milk)), "奶类 PQNC（供应商责任）",
    f"占供应商责任 {len(milk)/SP_JUL*100:.0f}%", CRIT, vsize=24)
kpi(s, 3.4, 1.4, 2.7, 1.05, "8", "效期标签缺失/模糊", "本月最高频子类", CRIT, vsize=24)
kpi(s, 6.25, 1.4, 2.7, 1.05, "7", "灌装量不足 / 液位偏低", "Cream-O-Land 集中", WARN, vsize=24)
kpi(s, 9.1, 1.4, 2.7, 1.05, str(len(milk_all)), "全月奶类相关合计",
    f"占全月 {len(milk_all)/N_JUL*100:.0f}%，其中共担集群 "
    f"{len([r for r in milk_all if resp(r)=='Joint'])} 起", INFO, vsize=24)
band(s, 0.55, 2.7, 12.2, "问题描述 / NC Description")
body(s, 0.55, 3.04, 12.2, 1.55,
     "效期标签（8 起）：7/8 字迹褪色；7/12、7/19、7/28 无效期；7/12 效期不清；7/23 单批 5 瓶无效期；7/24、7/26 标签未贴牢/脱落。\n"
     "灌装量不足 / 液位偏低（7 起）：7/2、7/3 ×3、7/6 液位明显低于常规；7/3、7/5 到货即已开封且内容物缺失。\n"
     "渗漏 / 破孔（4 起）：7/12 未开封瓶渗漏；7/16 瓶盖未完全盖合致脱脂奶溢出；7/27 瓶身侧面破孔；7/28 产品渗漏。\n"
     "变质 / 酸败（3 起）：7/5 色泽异常伴异味；7/19 颜色异常、酸味且液位偏低；7/21 椰奶开瓶即变质。", size=9.5)
band(s, 0.55, 4.75, 6.0, "问题分析 / Root Cause")
body(s, 0.55, 5.09, 6.0, 1.65,
     "效期喷码与标签粘附稳定性不足，为供应商灌装线工艺问题，非门店可闭环。\n"
     "灌装量不足与瓶盖压合不良指向同一条灌装线的稳定性。\n"
     "Cream-O-Land 为奶类单一主力供应商，风险高度集中——\n"
     f"7/22 起同一供应商再现 {len(JOINT)} 起变质集群（共担判责，见 Case B）。", size=9.5)
band(s, 6.75, 4.75, 6.0, "纠正措施 / Corrective Action（P1）")
body(s, 6.75, 5.09, 6.0, 1.65,
     "① 向 Cream-O-Land / SYSCO 提出灌装线与效期喷码专项整改，要求书面回复。\n"
     "② 门店收货增加「效期字迹清晰度 + 灌装液位 + 瓶盖压合」三项必检并拍照留证。\n"
     "③ 连续两月同类问题的批次纳入到货批次抽检加严。", size=9.5)

attach_note(s)   # CASE A: June carried on-site photos here

# ================================================================= 10 CASE B (new)
s = S(); title_bar(s, "典型案例 — Cream-O-Land Fat-Free Milk Spoilage Cluster",
                   f"脱脂奶变质集群 · Case B · 2026-07-22 ~ 07-28 · 判责：供应商·仓储共担", page())
j_stores = defaultdict(lambda: [0, 0.0, set()])
for r in JOINT:
    k = store(r); j_stores[k][0] += 1; j_stores[k][1] += val(r)
    j_stores[k][2].add(r["created_time"][:10])
j_dates = sorted({r["created_time"][:10] for r in JOINT})
kpi(s, 0.55, 1.4, 2.7, 1.05, str(len(JOINT)), "集群起数", "同一 SKU / 同一供应商", CRIT, vsize=24)
kpi(s, 3.4, 1.4, 2.7, 1.05, str(len(j_stores)), "涉及门店", f"{j_dates[0][5:]} ~ {j_dates[-1][5:]}", CRIT, vsize=24)
kpi(s, 6.25, 1.4, 2.7, 1.05, f"${JOINT_VAL:,.2f}", "损失货值", f"占全月 {JOINT_VAL/PJ_VAL*100:.0f}%（低值高频）", WARN, vsize=24)
kpi(s, 9.1, 1.4, 2.7, 1.05, "共担", "判责结果", "供应商 + 仓储物流", WARN, vsize=16)
table(s, 0.55, 2.65, 6.0, 2.55, ["门店 / Store", "Cases", "Value $", "日期"],
      [[k, v[0], f"${v[1]:,.2f}", " / ".join(d[5:] for d in sorted(v[2]))]
       for k, v in sorted(j_stores.items(), key=lambda kv: (-kv[1][0], kv[0]))] +
      [[("Total / 合计",{"bold":True}), (len(JOINT),{"bold":True}),
        (f"${JOINT_VAL:,.2f}",{"bold":True}), ""]],
      widths=[3.4,1.1,1.4,1.7], align_center=[1,2,3], fsize=8.5)
band(s, 6.75, 2.65, 6.0, "问题描述 / NC Description")
body(s, 6.75, 2.99, 6.0, 2.21,
     f"{len(JOINT)} 起工单全部指向同一 SKU：Cream-O-Land Fat-Free Milk 4/1 GAL CS（供应商 SYSCO），"
     f"涉及 {len(sorted({str(ENR[r['pqnc_no']]['batch']) for r in JOINT}))} 个批次"
     f"（{' / '.join(sorted({str(ENR[r['pqnc_no']]['batch']) for r in JOINT}))}）。\n"
     "门店描述高度一致：开瓶即有酸味 / 腐臭味、色泽异常、口感变酸，多数在效期内即已变质"
     "（如 7/28 一起效期标注 7/29）；1 起伴随瓶体渗漏。\n"
     f"7/22 单日集中上报 9 起，7/26、7/28 各再现 2 起，跨 {len(j_stores)} 家门店，非单店偶发。", size=9)
band(s, 0.55, 5.35, 6.0, "判责与分析 / Judgment & Root Cause")
body(s, 0.55, 5.69, 6.0, 1.32,
     "系统判责意见（2026-08-03）：多店反馈奶变酸，供应商与仓库双方自查\n"
     "均未发现异常 → 判为「供应商·仓储共担」，根因未落到单一环节。\n"
     "⚠ 检出滞后：客诉 7/8 起已连续反馈奶变质酸败（见第 33 页），\n"
     "内部首张变质 PQNC 为 7/19、集群 7/22——客户比工单早约两周。", size=9)
band(s, 6.75, 5.35, 6.0, "纠正措施 / Corrective Action（P1）")
body(s, 6.75, 5.69, 6.0, 1.32,
     "① 要求 SYSCO / Cream-O-Land 追溯上述批次原奶与灌装记录并书面回函。\n"
     "② 补齐链路温度证据：仓库出库温度、配送车温度记录、门店收货温度必检并留证。\n"
     "③ 涉事批次全量下架复检；8 月同 SKU 到货加严抽检（开瓶感官 + 效期）。\n"
     "④ 建立「同 SKU 跨店 3 起以上」自动预警，并把客诉纳入触发条件——\n"
     "   本月客诉信号早于 PQNC 约两周，可据此提前锁定批次。", size=9)

attach_note(s)   # CASE B: June carried on-site photos here

# ================================================================= 11 CASE C
s = S(); title_bar(s, "典型案例 — BOH Fridge Abnormal Defrost Cycle",
                   "后厨冰箱异常除霜致物料损失 · Case C · 2026-07-28 单日集群", page())
dfr = [r for r in PJ if "defrost" in (r["problem_description"] or "").lower()]
dfr_val = sum(val(r) for r in dfr)
kpi(s, 0.55, 1.4, 2.7, 1.05, str(len(dfr)), "同日上报起数", "2026-07-28 单日", CRIT, vsize=24)
kpi(s, 3.4, 1.4, 2.7, 1.05, f"${dfr_val:,.2f}", "损失货值", f"占全月 {dfr_val/PJ_VAL*100:.0f}%", CRIT, vsize=24)
kpi(s, 6.25, 1.4, 2.7, 1.05, "未明确", "当前判责", "非供应商 / 非仓储", WARN, vsize=16)
kpi(s, 9.1, 1.4, 2.7, 1.05, "9", "全月未明确判责", f"占全月 {9/N_JUL*100:.0f}%", WARN, vsize=24)
dfr_by = defaultdict(lambda: [0, 0.0])
for r in dfr:
    k = "Wholesome farms" if "wholesome" in (r["factory_name"] or "").lower() else "New York raw material warehouse"
    dfr_by[k][0] += 1; dfr_by[k][1] += val(r)
table(s, 0.55, 2.65, 6.0, 2.0, ["供应商 / 来源", "Cases", "Value $"],
      [[k, v[0], f"${v[1]:,.2f}"] for k, v in sorted(dfr_by.items(), key=lambda kv: -kv[1][0])] +
      [[("Total / 合计",{"bold":True}), (len(dfr),{"bold":True}), (f"${dfr_val:,.2f}",{"bold":True})]],
      widths=[5,1.6,2], align_center=[1,2], fsize=9)
band(s, 6.75, 2.65, 6.0, "问题描述 / NC Description")
body(s, 6.75, 2.99, 6.0, 1.66,
     f"{len(dfr)} 起工单描述一致：物料存放于后厨常温/冷藏冰箱，"
     "冰箱在非预期时段启动除霜循环（defrost cycle），导致温度失控、物料报废。\n"
     "全部由门店发起，判责结果为「未明确」，未归入供应商或仓储责任。", size=9.5)
band(s, 0.55, 4.8, 6.0, "问题分析 / Root Cause")
body(s, 0.55, 5.14, 6.0, 1.6,
     "属门店设备（冰箱控制板 / 除霜定时器）故障，非物料本身质量问题，\n"
     "因此现有 PQNC 判责体系无对应责任方，全部落入「未明确」。\n"
     f"本月「未明确」{PJ_RESP.get('Unknown/reject',0)} 起中 {len(dfr)} 起源于此，是该桶占比偏高的主因。", size=9.5)
band(s, 6.75, 4.8, 6.0, "纠正措施 / Corrective Action（P1）")
body(s, 6.75, 5.14, 6.0, 1.6,
     "① 立即排查涉事门店冰箱除霜定时器与控制板，必要时更换。\n"
     "② 冷藏/冷冻设备加装温度记录仪并纳入日检。\n"
     "③ 建议 PQNC 判责增设「门店设备」责任类型，避免设备类损失长期沉淀在「未明确」。", size=9.5)

attach_note(s)   # CASE C: June carried on-site photos here

# ================================================================= 12 CASE D
s = S(); title_bar(s, "典型案例 — Packaging Integrity",
                   "包装完整性（咖啡豆袋 / 罐体 / 杯具）· Case D", page())
band(s, 0.55, 1.35, 12.2, "问题明细 / Detail")
table(s, 0.55, 1.72, 12.2, 2.1,
      ["日期", "供应商 / 产品", "问题描述（原文）", "判责", "Value $"],
      [["07-05", "Sysco Corporation（罐装）", "Opened the can in the morning, found discolored object inside", "Supplier · 食安", "$2.26"],
       ["07-07", "Casa solana（罐装）", "Can dented along seam not suitable for consumption", "Supplier", "$2.26"],
       ["07-11", "S&D Coffee（咖啡豆）", "Broken bag when unboxing, beans scattered all over the inside of the box", "Supplier", "$13.61"],
       ["07-11", "S&D Coffee（咖啡豆）", "A bag was ripped completely opened with half the beans in the box", "Supplier", "$13.61"],
       ["07-28", "Luckin Medium Roast（咖啡豆）", "The bag itself is already ripped inside of the box", "Supplier", "$13.61"],
       ["07-09", "Pet bottle universal（杯具）", "No wellness label on the cup package", "Supplier", "$0.52"],
       ["07-22", "US-16oz ice cup NBJ（杯具）", "The cup has no brand logo, just an empty plain cup", "Supplier", "$0.07"]],
      widths=[1.2,3.2,6.2,1.8,1.2], align_center=[0,3,4], fsize=8.5)
band(s, 0.55, 4.0, 6.0, "问题分析 / Root Cause")
body(s, 0.55, 4.34, 6.0, 2.3,
     "咖啡豆袋（3 起）：开箱即见袋体破裂散豆，为封口强度或装箱挤压问题，\n"
     "S&D Coffee 与瑞幸自有中焙豆均出现，指向共用包材或装箱工艺。\n\n"
     "罐体（2 起）：7/5 罐内异物变色为本月唯一食安风险；7/7 罐体接缝凹陷。\n\n"
     "杯具（2 起）：无品牌标识 / 无 wellness 标签，属规格不符，\n"
     "与 6 月 HEC 4L 冷水壶无刻度线为同类「规格与采购标准不一致」问题。", size=9.5)
band(s, 6.75, 4.0, 6.0, "纠正措施 / Corrective Action（P2）")
body(s, 6.75, 4.34, 6.0, 2.3,
     "① 咖啡豆袋：要求供应商提升封口强度并复核装箱缓冲，连续跟踪 8 月到货。\n\n"
     "② 罐装产品：7/5 异物件已留样，要求 Sysco 追溯批次并书面回复根因。\n\n"
     "③ 杯具规格：入库增加规格抽检（品牌标识、标签、容量、材质），\n"
     "延续 6 月冷水壶抽检机制，避免规格不符流入门店。", size=9.5)

attach_note(s)   # CASE D: June carried on-site photos here

# ================================================================= 13 DIVIDER 03
s = S()
pic(s, "divider03", 0, 0, 7.73, 7.50)
section_head(s, "03", "Store Audit Performance\n门店稽核", 8.12, 3.10, 5.21, 1.72)

# ================================================================= 14 QA OVERALL
s = S(); title_bar(s, "7 月 QA 稽核整体表现",
                   f"QA 食安稽核 {QA_N} 次 / {QA_STORES} 家门店 · 分数为申诉后调整分", page())
kpi(s, 0.55, 1.45, 2.9, 1.5, str(QA_AVG), "QA 平均分（↑0.3 vs 6 月 91.8）", "申诉前原始分 76.7", GOOD)
kpi(s, 3.65, 1.45, 2.9, 1.5, f"{QA_STORES}/{NPRIM}", "QA 稽核覆盖门店", "3 家新店未纳入 QA", WARN)
kpi(s, 6.75, 1.45, 2.9, 1.5, str(QA_N), "QA 稽核次数", "6 月 16 次，+3", GOOD)
kpi(s, 9.85, 1.45, 2.9, 1.5, str(QA_S), f"QA 关键项 S（重点项 M {QA_M}）", "全部为水槽与管道", CRIT)
band(s, 0.55, 3.2, 12.2, "本月 QA 稽核状态总览")
CARD = [("QA 覆盖", f"{QA_STORES}/{NPRIM}", WARN), ("平均分", f"{QA_AVG}（+0.3 vs 6 月）", GOOD),
        ("关键项短板", "水槽与管道 / 气隙 Air Gap", CRIT), ("跨月异动", f"{len(REPEAT['repeat'])} 家连续两月复现", CRIT)]
for i, (lb, vl, col) in enumerate(CARD):
    x = 0.55 + i*3.09
    rect(s, x, 3.62, 2.92, 0.92, fill=FILL_L); rect(s, x, 3.62, 2.92, 0.06, fill=col)
    tb(s, x+0.14, 3.76, 2.64, 0.25, lb, size=9, color=INK2)
    tb(s, x+0.14, 4.02, 2.64, 0.42, vl, size=11.5, bold=True, color=col, spacing=0.95)
band(s, 0.55, 4.76, 12.2, "摘要")
body(s, 0.55, 5.1, 12.2, 1.6,
     f"7 月 QA 稽核 {QA_N} 次、覆盖 {QA_STORES} 家门店；申诉后平均分 {QA_AVG}，较 6 月 91.8 提升 0.3 分，达标率维持 100%。\n"
     f"⚠ 但申诉前原始均分仅 {QA_AVG_ORG}（6 月 78.4），QA 现场发现的问题不降反增——分数改善来自申诉获批而非现场改善。\n"
     f"最大系统性短板仍为设施 — Sinks and Pipes / air gap：QA 主巡检 {QA_S} 个关键项全部为此类，且全部申诉获批扣分归零；\n"
     f"{len(REPEAT['repeat'])} 家门店连续两月复现同类问题，6 月 P0（BD 整改清单）未见成效，需升级为 BD 专项工程。", size=10)

# ================================================================= 15 THREE TYPES
s = S(); title_bar(s, "三类稽核体系对比", "门店自检 + QA 稽核 + 区经检查 · 本月三类齐全且全部活跃", page())
s71 = PACK["s7_1"]
table(s, 1.2, 1.5, 10.9, 2.0,
      ["稽核类型", "次数", "覆盖门店", "平均分", "关键项 S", "重点项 M"],
      [["门店自检", s71["门店自检"]["count"], f"{s71['门店自检']['stores']}/{NPRIM}", s71["门店自检"]["avg"], s71["门店自检"]["S"], s71["门店自检"]["M"]],
       ["QA 稽核", QA_N, f"{QA_STORES}/{NPRIM}", QA_AVG, s71["QA审计"]["S"], s71["QA审计"]["M"]],
       ["区经检查", s71["区经检查"]["count"], f"{s71['区经检查']['stores']}/{NPRIM}", s71["区经检查"]["avg"], s71["区经检查"]["S"], s71["区经检查"]["M"]]],
      widths=[2.4,1.2,1.4,1.3,1.3,1.3], align_center=[1,2,3,4,5], fsize=10.5, hsize=10)
band(s, 1.2, 3.75, 10.9, "体系表现")
body(s, 1.2, 4.09, 10.9, 2.55,
     f"7 月共 {TC['total']} 次有效巡检（自检 {TC['门店自检']} / QA {TC['QA审计']} / 区经 {TC['区经检查']}），较 6 月 85 次 +9 次，为有记录以来最高；\n"
     f"{NPRIM} 家在营门店 100% 覆盖，含 3 家新开业门店首检落地。\n\n"
     f"QA 稽核平均分最高（{QA_AVG}），{s71['QA审计']['S']} 个 S 项全部集中在 Sinks and Pipes / air gap，说明设施类问题仍未根治。\n\n"
     f"⚠ 尺度分化加剧：区经均分由 6 月 84.3 降至 {s71['区经检查']['avg']}，QA（逐店口径）由 6 月 91.8 升至 {QA_AVG}，"
     f"两者差距 {QA_AVG - s71['区经检查']['avg']:.1f} 分为历史最大；{len(PACK['s4_4'])} 家门店出现 ≥20 分跨类型评分背离"
     f"（6 月仅 4 家），需开展区经 / QA 尺度对齐校准。\n"
     f"※ 表中 QA 平均分为「每店最近一次 QA 稽核」口径（沿用 6 月 PPT 惯例）；自检 / 区经为全部单据均分。",
     size=10)

# ================================================================= 16 QA RANKING
s = S(); title_bar(s, "QA 门店得分排名", f"{QA_STORES} 家自营门店（申诉后分）· 80 分为达标线", page())
rank = sorted(QA_LATEST.items(), key=lambda kv: -kv[1]["adjusted_total_score"])
chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, 0.5, 1.45, 12.3, 3.9,
      [sn(c) for c, _ in rank], [("QA 得分", [r["adjusted_total_score"] for _, r in rank])],
      [C1], legend=False, lsize=8)
band(s, 0.5, 5.5, 12.3, "分档")
body(s, 0.5, 5.84, 12.3, 1.05,
     f"A+ 档 ≥94（{len([1 for _,r in rank if r['adjusted_total_score']>=94])} 家）："
     + "、".join(f"{sn(c)} {r['adjusted_total_score']}" for c, r in rank if r["adjusted_total_score"] >= 94) + "\n"
     f"全部 QA 稽核门店均 ≥84 分，无 <80 门店。最低分为 {sn(rank[-1][0])} {rank[-1][1]['adjusted_total_score']}、"
     f"{sn(rank[-2][0])} {rank[-2][1]['adjusted_total_score']}。", size=9)

# ================================================================= 17 QA GRADING
s = S(); title_bar(s, "QA 稽核分级结果",
                   f"自营 {QA_STORES} 家 · 申诉后调整分（A+ ≥94 / A 87–93 / B 80–86 / C <80）", page())
n = QA_STORES
chart(s, XL_CHART_TYPE.BAR_STACKED, 0.55, 1.5, 6.7, 2.3, ["QA 稽核"],
      [("A+ 级", [GA["A+"]]), ("A 级", [GA["A"]]), ("B 级", [GA["B"]]), ("C 级", [GA["C"]])],
      SEQ[:4], lsize=9, num_fmt="0")
kpi(s, 7.5, 1.5, 2.5, 1.1, f"{(GA['A+']+GA['A']+GA['B'])/n*100:.0f}%", "达标率（A+/A/B，≥80 分）", None, GOOD, vsize=18)
kpi(s, 10.25, 1.5, 2.5, 1.1, f"{(GA['A+']+GA['A'])/n*100:.1f}%", "A 级以上占比（A+/A）", None, GOOD, vsize=18)
table(s, 7.5, 2.85, 5.25, 0.95, ["类型", "稽核店数", "A+ 级", "A 级", "B 级", "C 级"],
      [["QA 稽核", n, f"{GA['A+']/n*100:.1f}%", f"{GA['A']/n*100:.1f}%", f"{GA['B']/n*100:.1f}%", f"{GA['C']/n*100:.1f}%"],
       ["家数", n, GA["A+"], GA["A"], GA["B"], GA["C"]]],
      widths=[1.6,1.3,1.1,1.1,1.1,1.1], align_center=[1,2,3,4,5], fsize=9)
band(s, 0.55, 4.15, 12.2, "总结")
body(s, 0.55, 4.49, 12.2, 2.2,
     f"QA 达标率 100%（{n}/{n}），A 级以上占比 {(GA['A+']+GA['A'])/n*100:.1f}%（{GA['A+']+GA['A']}/{n}），其中 A+ 门店 {GA['A+']} 家（6 月 6 家）。\n\n"
     f"本月无 C 级不合格门店。北美门店全部为自营。\n\n"
     f"⚠ 需注意：上述分级基于申诉后调整分。若按申诉前原始分分级，C 级门店达 {GO['C']} 家（{GO['C']/n*100:.1f}%），"
     f"达标率仅 {(GO['A+']+GO['A']+GO['B'])/n*100:.1f}%，详见下页。", size=10)

# ================================================================= 18 PRE VS POST APPEAL
s = S(); title_bar(s, "申诉前 vs 申诉后 分数对比",
                   "QA 稽核 · 申诉机制对评分结果的影响（申诉前 = 原始分，申诉后 = 调整分）", page())
kpi(s, 0.55, 1.45, 2.9, 1.1, str(QA_AVG_ORG), "QA 平均分（申诉前）", None, CRIT, vsize=24)
kpi(s, 3.65, 1.45, 2.9, 1.1, str(QA_AVG), f"QA 平均分（申诉后，+{QA_AVG-QA_AVG_ORG:.1f}）", None, GOOD, vsize=24)
tb(s, 6.9, 1.4, 5.9, 0.25, "分级结构：申诉前 vs 申诉后（家数）", size=9, bold=True, color=INK2)
chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, 6.75, 1.68, 6.0, 2.5, ["A+", "A", "B", "C"],
      [("申诉前", [GO["A+"], GO["A"], GO["B"], GO["C"]]),
       ("申诉后", [GA["A+"], GA["A"], GA["B"], GA["C"]])], [C2, C1], lsize=8.5)
band(s, 0.55, 2.72, 6.0, "申诉调整门店（申诉前 → 申诉后）")
adj_rows = sorted([(c, r["original_total_score"], r["adjusted_total_score"])
                   for c, r in QA_LATEST.items() if r["adjusted_total_score"] != r["original_total_score"]],
                  key=lambda x: -(x[2]-x[1]))
table(s, 0.55, 3.09, 6.0, 3.0, ["门店", "申诉前", "申诉后", "变动"],
      [[sn(c), o, a, (f"+{a-o}", {"bold": a-o >= 20, "color": CRIT if a-o >= 20 else INK})]
       for c, o, a in adj_rows],
      widths=[3.2,1.2,1.2,1.2], align_center=[1,2,3], fsize=8.5)
body(s, 6.75, 4.35, 6.0, 2.3,
     f"申诉调整使 QA 平均分由 {QA_AVG_ORG} 提升至 {QA_AVG}（+{QA_AVG-QA_AVG_ORG:.1f}），"
     f"达标率由 {(GO['A+']+GO['A']+GO['B'])/n*100:.1f}% 升至 100%，A 级以上由 {(GO['A+']+GO['A'])/n*100:.1f}% 升至 {(GA['A+']+GA['A'])/n*100:.1f}%。\n\n"
     f"本月共 {len(adj_rows)} 家门店经申诉调整，全部为系统获批（6 月为 7 起系统获批 + 4 起 QA 手动更新）。\n\n"
     f"其中 10 起为 air gap / 管道 S 项撤销，单笔提分 +25 至 +31（S 项撤销同时解除 −20 惩罚项）。\n\n"
     f"⚠ 连续两月已决申诉 0 驳回（累计 23 获批）。建议引入第二审核人、S 项申诉强制附整改证据、获批后 14 天内实物复检。", size=9)

# ================================================================= 19 JUL VS JUN
s = S(); title_bar(s, "QA 稽核 7 月 vs 6 月 对比", None, page())
kpi(s, 0.55, 1.5, 3.9, 1.1, f"91.8 → {QA_AVG}", "QA 平均分 6月→7月（+0.3）", None, GOOD, vsize=19)
kpi(s, 4.65, 1.5, 3.9, 1.1, f"16 → {QA_N}", "QA 稽核次数 6月→7月", None, GOOD, vsize=19)
kpi(s, 8.75, 1.5, 4.0, 1.1, f"78.4 → {QA_AVG_ORG}", "申诉前均分 6月→7月（−1.7）", None, CRIT, vsize=19)
tb(s, 0.55, 2.85, 6.5, 0.25, "分级结构对比（家数）", size=9, bold=True, color=INK2)
chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, 0.5, 3.1, 7.0, 2.7, ["A+", "A", "B", "C"],
      [("6 月", [6, 8, 2, 0]), ("7 月", [GA["A+"], GA["A"], GA["B"], GA["C"]])], [C2, C1], lsize=8.5)
band(s, 7.8, 3.1, 5.0, "要点")
body(s, 7.8, 3.44, 5.0, 2.36,
     f"达标结构维持高位：\n\n"
     f"· A+ 门店由 6 家增至 {GA['A+']} 家；C 级维持 0 家。\n\n"
     f"· QA 覆盖由 16 家增至 {QA_STORES} 家（+2），3 家新店仍未纳入。\n\n"
     f"· S 项由 8 项增至 {QA_S} 项，仍全部集中在设施 / 管道。\n\n"
     f"⚠ 申诉前均分由 78.4 降至 {QA_AVG_ORG}，现场问题实际增多。", size=9.5)

# ================================================================= 20 YTD
s = S(); title_bar(s, "2026 年 YTD 概览（1–7 月）", None, page())
ytd_self = 7+5+13+32+49+51+TC["门店自检"]; ytd_qa = 5+2+1+12+16+16+TC["QA审计"]
ytd_area = 4+0+0+14+21+18+TC["区经检查"]; ytd_tot = ytd_self+ytd_qa+ytd_area
ytd_s = 84 + FM["S"]
kpi(s, 0.55, 1.5, 3.0, 1.1, str(ytd_tot), "YTD 稽核总次数（1–7 月）", None, INFO, vsize=24)
kpi(s, 3.75, 1.5, 3.0, 1.1, str(ytd_qa), "YTD QA 稽核次数", None, INFO, vsize=24)
kpi(s, 6.95, 1.5, 3.0, 1.1, str(ytd_s), "YTD 关键项 S 累计", None, CRIT, vsize=24)
kpi(s, 10.15, 1.5, 2.6, 1.1, f"{NPRIM} 家", "YTD 已覆盖自营门店", None, GOOD, vsize=24)
tb(s, 0.55, 2.85, 6.5, 0.25, "稽核量月度趋势（按类型，次）", size=9, bold=True, color=INK2)
chart(s, XL_CHART_TYPE.COLUMN_STACKED, 0.5, 3.1, 7.4, 2.8,
      ["1月","2月","3月","4月","5月","6月","7月"],
      [("门店自检", [7,5,13,32,49,51,TC["门店自检"]]),
       ("QA 稽核", [5,2,1,12,16,16,TC["QA审计"]]),
       ("区经检查", [4,0,0,14,21,18,TC["区经检查"]])], [C1, C2, C3], lsize=8)
band(s, 8.2, 3.1, 4.6, "YTD 要点")
body(s, 8.2, 3.44, 4.6, 2.46,
     f"1–7 月共 {ytd_tot} 次有效巡检\n（自检 {ytd_self} / QA {ytd_qa} / 区经 {ytd_area}）。\n\n"
     f"7 月巡检体系创历史新高：\n自检 {TC['门店自检']}、QA {TC['QA审计']}、区经 {TC['区经检查']}。\n\n"
     f"门店数由 1 月 13 家扩至 7 月 {NPRIM} 家，\n覆盖率连续 4 个月保持 100%。", size=9.5)

# ================================================================= 21 STABILITY
s = S(); title_bar(s, "门店稳定性分析（6 月 → 7 月 主巡检）", "逐店对比 · 识别下滑 / 稳定 / 改善", page())
imp = PACK["s1_3"]["improvers"]; dec = PACK["s1_3"]["decliners"]; flat = PACK["s1_3"]["flat"]
big_imp = [x for x in imp if x["delta"] >= 5]; small = [x for x in imp if x["delta"] < 5] + flat + [x for x in dec if x["delta"] > -5]
big_dec = [x for x in dec if x["delta"] <= -5]
COLS = [("下滑（≥5 分）", big_dec, CRIT), ("稳定（±4 分内）", small, C1),
        ("改善（≥5 分）", big_imp, GOOD), ("新店首检", None, WARN)]
for i, (hd, items, col) in enumerate(COLS):
    x = 0.55 + i*3.09
    rect(s, x, 1.5, 2.92, 0.42, fill=col)
    tb(s, x+0.12, 1.59, 2.68, 0.25, hd, size=10, bold=True, color=WHITE)
    rect(s, x, 1.92, 2.92, 3.35, fill=SURF, line=RULE)
    if items is None:
        txt = ("48th & 3rd 60（区经）\nGrand Central Terminal 96（区经）\n128 W 32nd St 64（区经）\n\n"
               "三家新店本月均未纳入 QA 稽核，首检全部由区经完成。\n\n"
               "US00021 于 7/16 开业、7/28 完成首检。")
    elif not items:
        txt = "无"
    else:
        txt = "\n".join(f"{sn(v['store'])} {v['jun']}→{v['jul']}（{v.get('delta', v['jul']-v['jun']):+d}）"
                        for v in items)
    tb(s, x+0.14, 2.06, 2.64, 3.05, txt, size=8.5, color=INK, spacing=1.3)
band(s, 0.55, 5.42, 12.2, "总结")
body(s, 0.55, 5.76, 12.2, 1.1,
     f"重点观察：本月无 ≥5 分下滑门店，最大跌幅仅 −3（21st & 3rd 97→94），门店表现整体趋稳。\n"
     f"52nd & Madison 由 6 月最低分（64）回升至 90（+26），54th & 8th 71→84（+13）结束连续两月 <80；两家 6 月 <80 门店全部脱困。"
     f"新增关注对象转为 3 家新店中的 48th & 3rd（60）与 128 W 32nd St（64）。", size=10)

# ================================================================= 22 SEVERITY TREND
s = S(); title_bar(s, "稽核问题点分析", "自营门店关键项 / 重点项 / 一般项 / 轻微项趋势 · 6 月 → 7 月（主巡检口径，店均）", page())
jun_n, jul_n = 18, NPRIM
chart(s, XL_CHART_TYPE.COLUMN_STACKED, 0.55, 1.55, 6.5, 4.5, ["6 月", "7 月"],
      [("店均轻微项 L", [round(16/jun_n,2), round(PM["L"]/jul_n,2)]),
       ("店均一般项 G", [round(58/jun_n,2), round(PM["G"]/jul_n,2)]),
       ("店均重点项 M", [round(6/jun_n,2), round(PM["M"]/jul_n,2)]),
       ("店均关键项 S", [round(8/jun_n,2), round(PM["S"]/jul_n,2)])],
      SEQ[:4], lsize=8, num_fmt="0.00")
band(s, 7.35, 1.55, 5.4, "关键变化")
body(s, 7.35, 1.89, 5.4, 4.16,
     f"店均关键项 S 上升：0.44 → {PM['S']/jul_n:.2f}\n"
     f"7 月主巡检 {PM['S']} 个 S 项中 10 个为 Sinks and Pipes / air gap，"
     f"另 3 个来自新店（交叉污染 2、洗手规范 1），设施类短板仍未根治。\n\n"
     f"店均重点项 M 持平：0.33 → {PM['M']/jul_n:.2f}\n"
     f"7 起 M 项中 5 起为开封后效期标签缺失，集中度极高，属可运营闭环项。\n\n"
     f"店均一般项 G 上升：3.22 → {PM['G']/jul_n:.2f}\n"
     f"主要为清洁卫生类（49 项），冰箱顶积尘、器具残留、消毒残留为最高频描述。\n\n"
     f"店均轻微项 L 下降：0.89 → {PM['L']/jul_n:.2f}\n\n"
     f"※ 口径：均以当月主巡检门店数为分母（6 月 18 家 / 7 月 {jul_n} 家）。", size=9)

# ================================================================= 23 MODULE RISK
s = S(); title_bar(s, "稽核模块风险分析", f"主巡检 8 模块扣分排名（{NPRIM} 家门店）", page())
mods = sorted([m for m in MA if MA[m]["problems"] > 0], key=lambda m: MA[m]["deduction"])
chart(s, XL_CHART_TYPE.BAR_CLUSTERED, 0.5, 1.5, 7.4, 4.9,
      [m for m in mods][::-1], [("扣分", [abs(MA[m]["deduction"]) for m in mods][::-1])],
      [C2], legend=False, lsize=8.5)
band(s, 8.2, 1.5, 4.55, "风险分层")
body(s, 8.2, 1.84, 4.55, 4.56,
     f"系统性问题（≥50% 门店）\n"
     f"清洁卫生（100%，{MA['清洁卫生']['deduction']}）\n"
     f"设施（{MA['设施']['coverage']}%，{MA['设施']['deduction']}，含 10 个关键项）\n\n"
     f"中等覆盖面（30–49%）\n本月无\n\n"
     f"低覆盖面问题（<30%）\n"
     f"温控有效期 {MA['温控有效期']['deduction']}（全部为 M 项 · 效期标签）\n"
     f"过程控制 {MA['过程控制']['deduction']}（含 2 个新店关键项）\n"
     f"虫害防控 {MA['虫害防控']['deduction']}（本月新增进入主巡检）\n"
     f"职业安全 {MA['职业安全']['deduction']}、员工健康卫生 {MA['员工健康卫生']['deduction']}、设备维护 {MA['设备维护']['deduction']}\n\n"
     f"证照 / 供应链 本月主巡检无扣分。", size=9.5)

# ================================================================= 24 S & M DETAIL
s = S(); title_bar(s, "稽核关键项 S & 重点项 M", f"主巡检关键项 {PM['S']} 项 / 重点项 {PM['M']} 项", page())
band(s, 0.55, 1.4, 6.0, f"关键项（S）明细 — 共 {PM['S']} 项")
s32 = PACK["s3_2"]
SUBCN = {"Sinks and Pipes":"水槽与管道","Cross-Contamination":"交叉污染","Handwashing Standards":"洗手规范"}
rows_s = [[sn(x["store"]), SUBCN.get(x["sub_item"], x["sub_item"]),
           (x["description"] or "").replace("\n"," ")[:44], x["deduction"]] for x in s32[:9]]
table(s, 0.55, 1.77, 6.0, 3.3, ["门店", "子类", "问题描述", "扣分"], rows_s,
      widths=[2.0,1.5,4.4,0.8], align_center=[3], fsize=7.5)
tb(s, 0.55, 5.12, 6.0, 0.22, f"※ 列示前 9 项，其余 {PM['S']-9} 项见月度分析报告 docx §3.2",
   size=7.5, color=MUTED)
band(s, 6.75, 1.4, 6.0, f"重点项（M）明细 — 共 {PM['M']} 项")
rows_m = [[sn(x["store"]), x["module"], (x["description"] or "").replace("\n"," ")[:40], x["deduction"]]
          for x in PACK["s3_4"]]
table(s, 6.75, 1.77, 6.0, 3.5, ["门店", "模块", "问题描述", "扣分"], rows_m,
      widths=[2.0,1.5,4.4,0.8], align_center=[3], fsize=7.5)
body(s, 0.55, 5.45, 12.2, 1.4,
     f"关键项 S 高度集中：{PM['S']} 个主巡检 S 项中 {len(PRIM_PIPE)} 个为「水槽与管道 / 气隙」"
     f"（{'、'.join(sn(c) for c in PRIM_PIPE)}）；另 3 个来自新店区经检查。\n"
     f"⚠ 10 个 air gap 关键项已全部经申诉获批、扣分归零，但物理问题未整改（详见下页跨月台账）。\n"
     f"重点项 M：效期管理 5 项、清洁与消毒 1 项、虫害防控 1 项（15th & 3rd 发现活蟑螂，需 7 天内闭环）。", size=9.5)

# ================================================================= 25 AIR GAP CASE
s = S(); title_bar(s, "典型案例 — 水槽与管道 / 气隙（关键项 S）",
                   f"设施模块 · 主巡检 10 个关键项属此类 · 全月 17 起、跨 11 家门店 · 连续第三个月未根治", page())
kpi(s, 0.55, 1.4, 2.9, 1.0, "10 项", "主巡检关键项（水槽与管道）", None, CRIT, vsize=19)
kpi(s, 3.65, 1.4, 2.9, 1.0, "17 起", "全月（含自检 / 区经）", "6 月 13 起，+4", CRIT, vsize=19)
kpi(s, 6.75, 1.4, 2.9, 1.0, "11 家", "涉及门店（全月）", "6 月 10 家", CRIT, vsize=19)
kpi(s, 9.85, 1.4, 2.9, 1.0, "连续 3 月", "相同问题未解决", f"{len(REPEAT['repeat'])} 家两月复现", CRIT, vsize=16)
band(s, 0.55, 2.6, 12.2, "问题描述")
body(s, 0.55, 2.94, 12.2, 1.4,
     "54th & 8th：milk machine 后方管道装置不符标准。　102 Fulton：冰机与滴滤机下方地漏 air gap 不合规。\n"
     "154 Bleecker：洗手池 / 三槽池排水持续外溢，滴滤机下方地漏 backflow。　40th & 10th：管道距滤芯过近、包覆层脱落。\n"
     "21st & 3rd：油脂阱管道触碰滤芯、咖啡机下方 air gap 不合规。　221 Grand：滴滤机下方管道触碰滤芯。\n"
     "33rd & 10th / 15th & 3rd / 52nd & Madison：air gap 不符标准（含三槽池）。　29th & 3rd：管道包覆异常。", size=9.5)
band(s, 0.55, 4.5, 6.0, "问题分析")
body(s, 0.55, 4.84, 6.0, 1.95,
     "属设施 / 管道结构性问题，非门店日常操作可闭环，需 BD / 营建介入。\n\n"
     "4–6 月已连续识别同类 air gap 短板，7 月仍跨 11 家门店出现且起数上升，\n"
     "说明「整改工单 + 申诉调分」的现行机制未能推动实物整改。", size=9.5)
band(s, 6.75, 4.5, 6.0, "整改建议（P0）")
body(s, 6.75, 4.84, 6.0, 1.95,
     "① 由「整改工单」升级为「BD 专项工程」，逐店出具管道 / air gap 改造方案与完工验收单。\n"
     "② QA 复检以实物验收为准，不接受仅申诉调分。\n"
     "③ 对连续两月复现门店设定 8 月底前完工时限，未闭环纳入门店考核。", size=9.5)

# ================================================================= 26 CROSS-MONTH LEDGER
s = S(); title_bar(s, "Air Gap 跨月追踪台账（6 月 → 7 月）",
                   "落实 6 月报告 P0 要求建立的跨月关键项台账 · 逐店记录复现状态", page())
jun_set, jul_set = set(REPEAT["jun_stores"]), set(REPEAT["jul_stores"])
led = []
for c in sorted(jun_set | jul_set):
    inj = "✓" if c in jun_set else "—"; inl = "✓" if c in jul_set else "—"
    if c in jun_set and c in jul_set: st, col = "连续两月复现（未根治）", CRIT
    elif c in jul_set:                st, col = "7 月新增", WARN
    else:                             st, col = "7 月已消除", GOOD
    led.append([sn(c), c, inj, inl, (st, {"color": col, "bold": True})])
table(s, 0.55, 1.45, 6.1, 4.3, ["门店", "编号", "6 月", "7 月", "跨月状态"], led,
      widths=[2.3,1.4,0.8,0.8,2.4], align_center=[1,2,3,4], fsize=8.5)
kpi(s, 6.9, 1.45, 1.85, 1.0, str(len(REPEAT["repeat"])), "连续两月复现", None, CRIT, vsize=24)
kpi(s, 8.9, 1.45, 1.85, 1.0, str(len(REPEAT["new"])), "7 月新增", None, WARN, vsize=24)
kpi(s, 10.9, 1.45, 1.85, 1.0, str(len(REPEAT["cleared"])), "7 月已消除", None, GOOD, vsize=24)
band(s, 6.9, 2.7, 5.85, "台账结论")
body(s, 6.9, 3.04, 5.85, 2.7,
     f"6 月 {len(jun_set)} 家 → 7 月 {len(jul_set)} 家；连续两月复现 {len(REPEAT['repeat'])} 家、"
     f"新增 {len(REPEAT['new'])} 家（{'、'.join(sn(c) for c in REPEAT['new'])}）、"
     f"已消除 {len(REPEAT['cleared'])} 家（{'、'.join(sn(c) for c in REPEAT['cleared'])}）。\n\n"
     f"结论：6 月 P0（48 小时闭环 + BD 整改清单）执行未见成效——"
     f"{len(REPEAT['repeat'])}/{len(jun_set)} 的 6 月问题门店在 7 月再次命中同类关键项，"
     f"且全月起数由 13 升至 17。\n\n"
     f"这些门店的 QA 分数已通过申诉恢复至 90+，但物理管道问题仍在，"
     f"形成「分数已修复、风险未消除」的错配，是本月最需管理层介入的事项。", size=9.5)

# ================================================================= 27 M CASE
s = S(); title_bar(s, "典型案例 — 效期标签与消毒管理（重点项 M）",
                   "开封后效期标签缺失为本月最集中重点项（5/7）", page())
band(s, 0.55, 1.4, 12.2, "问题描述")
body(s, 0.55, 1.74, 12.2, 1.5,
     "100 Maiden Ln：产品缺少效期日期。　54th & 8th：缺少效期标签；消毒液读数低于标准。\n"
     "102 Fulton：糖浆瓶缺少效期日期。　16th & 6th：糖浆瓶缺少效期标签。　128 W 32nd St（新店）：无标签。\n"
     "15th & 3rd：楼下发现活蟑螂爬行，捕虫器已满（虫害防控 M 项，需 7 天内闭环并复检）。", size=9.5)
band(s, 0.55, 3.35, 6.0, "问题分析")
body(s, 0.55, 3.69, 6.0, 2.9,
     "效期标签缺失 5 起中，3 起为糖浆瓶——集中于同一物料类型，\n"
     "指向吧台开封后贴标环节的执行断点，而非普遍性 SOP 缺失。\n\n"
     "该类问题属门店日常 SOP 执行与班次检查不到位，可完全运营闭环。\n\n"
     "与供应链侧 8 起「奶类效期标签缺失」形成呼应：\n"
     "门店端与供应商端同时存在效期标识管理薄弱。", size=9.5)
band(s, 6.75, 3.35, 6.0, "整改建议（P1）")
body(s, 6.75, 3.69, 6.0, 2.9,
     "① 开封后效期标签纳入每日开店清单必检项，糖浆瓶为重点。\n\n"
     "② 建议标签打印机点位前置到吧台，减少贴标动线断点。\n\n"
     "③ 消毒液浓度（PPM）每班次校准并留痕。\n\n"
     "④ 15th & 3rd 立即启动灭治并于 7 日内复检；全店排查虫控服务报告归档完整性"
     "（23rd & 8th 6 月服务报告缺失）。", size=9.5)

# stores where Orkin logged any activity — needed one page early for the cross-ref
p_by_store_pre = {o["store_name"] for o in PEST["observations"]
                  if "Activity" in o["observation"] and "No Activity" not in o["observation"]}

# ================================================================= 28 PEST (NEW)
s = S(); title_bar(s, "新增风险面 — 虫害防控（Pest Control）",
                   "本月首次进入主巡检 · 6 月主巡检 0 项 → 7 月 4 项 / 3 家门店", page())
kpi(s, 0.55, 1.45, 2.9, 1.1, "4 项", "主巡检虫害防控扣分项", "6 月 0 项", CRIT, vsize=24)
kpi(s, 3.65, 1.45, 2.9, 1.1, "3 家", "涉及门店", f"覆盖率 {MA['虫害防控']['coverage']}%", WARN, vsize=24)
kpi(s, 6.75, 1.45, 2.9, 1.1, "1", "重点项 M（活蟑螂）", "15th & 3rd", CRIT, vsize=24)
kpi(s, 9.85, 1.45, 2.9, 1.1, f"{MA['虫害防控']['deduction']}", "主巡检扣分", "全月 15 项", WARN, vsize=24)
band(s, 0.55, 2.75, 12.2, "问题明细")
table(s, 0.55, 3.12, 12.2, 1.7, ["门店", "子项", "严重度", "问题描述（原文）", "扣分"],
      [["15th & 3rd", "No Sign of Insect Pests", "M", "Live cockroach found running across floor. (Downstairs) / Traps filled with bugs.", "-5"],
       ["21st & 3rd", "No Sign of Insect Pests", "G", "Flies seen on drainage system.", "-2"],
       ["21st & 3rd", "Prevent pests from outside", "G", "Air curtain out of commission.", "-2"],
       ["23rd & 8th", "Pest control devices", "G", "Service report for June missing.", "-2"]],
      widths=[1.6,2.4,0.9,6.3,0.8], align_center=[2,4], fsize=8.5)
band(s, 0.55, 5.0, 6.0, "问题分析")
body(s, 0.55, 5.34, 6.0, 1.55,
     "三家门店问题互不相同但同指虫控体系薄弱：\n"
     "实体虫害（活蟑螂 / 蝇）、物理屏障失效（风幕机停用）、\n"
     "以及服务记录缺失（6 月虫控报告未归档）。\n"
     "该模块 6 月主巡检为 0，属新出现风险面，需即时干预。\n\n"
     f"⚠ 但主巡检并非全貌：第三方服务商 Orkin 7 月在 {len(p_by_store_pre)} 家门店\n"
     "发现虫害痕迹，与本页仅重合 1 家 —— 详见下页。", size=9.5)
band(s, 6.75, 5.0, 6.0, "整改建议（P1）")
body(s, 6.75, 5.34, 6.0, 1.55,
     "① 15th & 3rd 立即启动灭治并 7 日内复检，捕虫器更换周期改为每周。\n"
     "② 21st & 3rd 风幕机立即报修；排水系统清理消杀。\n"
     "③ 全店排查虫控服务报告归档完整性，服务改为月度双向签认——\n"
     "   7 月 21 家门店服务报告已全部到齐（见下页），归档缺口已补上。", size=9.5)

# ================================================== 29 PEST — VENDOR SERVICE VIEW (NEW)
s = S(); title_bar(s, "虫害防控 — 第三方服务商视角（Orkin）",
                   f"21 家门店月度例行消杀 · {PD['date_range'][0]} ~ {PD['date_range'][1]} · "
                   f"数据源：July Service Report（逐店 PDF 服务报告）", page())
p_act = [o for o in PEST["observations"]
         if "Activity" in o["observation"] and "No Activity" not in o["observation"]]
p_live = [o for o in p_act if "Live" in o["observation"]]
p_by_store = defaultdict(list)
for o in p_act:
    p_by_store[o["store_name"]].append(o)
QA_PEST = {"15th & 3rd", "21st & 3rd", "23rd & 8th"}   # 主巡检点名门店（见上页）
both = QA_PEST & set(p_by_store)
union = QA_PEST | set(p_by_store)

kpi(s, 0.55, 1.45, 2.9, 1.1, str(PD["visits"]), "服务门店 / 覆盖", "每店 1 次月度例行", GOOD)
kpi(s, 3.65, 1.45, 2.9, 1.1, str(PD["live"]), "活体发现 / Live activity",
    f"{PD['stores_with_live']} 家门店", CRIT)
kpi(s, 6.75, 1.45, 2.9, 1.1, str(len(p_by_store)), "有虫害痕迹门店",
    "活体 + 死体，主巡检仅点名 3 家", CRIT)
kpi(s, 9.85, 1.45, 2.9, 1.1, f"${PD['spend_total']:,.0f}", "本月服务费", "21 次 × $97.99", INFO)
band(s, 0.55, 2.72, 7.4, "虫害痕迹明细 / Activity found by Orkin")
# the 108th & Broadway report repeats one block verbatim — collapse for display,
# the raw duplicate stays in the data pack (flagged in july2026_pest_validation.txt)
p_rows, _seen = [], set()
for o in sorted(p_act, key=lambda o: (0 if "Live" in o["observation"] else 1, o["store_name"])):
    key = (o["store_name"], o["observation"], o["pest_type"], o["location"])
    if key in _seen:
        continue
    _seen.add(key)
    p_rows.append([o["store_name"], "活体" if "Live" in o["observation"] else "死体",
                   o["pest_type"] or "—", (o["location"] or "—")[:26]])
table(s, 0.55, 3.09, 7.4, 2.35, ["门店", "活体/死体", "虫种", "位置"], p_rows,
      widths=[2.0,1.1,1.9,2.6], align_center=[1], fsize=8.5)
band(s, 8.2, 2.72, 4.55, "与主巡检的交叉核对")
body(s, 8.2, 3.09, 4.55, 2.35,
     f"主巡检点名 {len(QA_PEST)} 家（上页）：{'、'.join(sorted(QA_PEST))}\n\n"
     f"Orkin 发现痕迹 {len(p_by_store)} 家。\n\n"
     f"「两者仅重合 {len(both)} 家（{'、'.join(sorted(both)) or '无'}）」，"
     f"合并口径下本月有虫害信号的门店为 「{len(union)} 家」——\n"
     f"是主巡检单一视角的 {len(union)/len(QA_PEST):.1f} 倍。", size=9)
band(s, 0.55, 5.58, 6.0, "问题分析")
body(s, 0.55, 5.92, 6.0, 1.0,
     f"· 活体 {PD['live']} 起中 2 起在地漏/排水（33rd & 10th、52nd & Madison），\n"
     f"  与主巡检 21st & 3rd「蝇集中于排水系统」同源；\n"
     f"· 死体以地下室/储藏区美洲蟑螂为主（102 Fulton、108th & Broadway、15th & 3rd）；\n"
     f"· 6 月遗留 {PD['open_actions_unresolved']} 项预防处理仍未闭环（含 21st & 3rd）。", size=8.5)
band(s, 6.75, 5.58, 6.0, "整改建议（P1）")
body(s, 6.75, 5.92, 6.0, 1.0,
     "① 排水/地漏专项：清理消杀 + 加装防虫网，按周复检至连续两次无蝇。\n"
     "② 地下室/储藏区：美洲蟑螂重点布控，捕虫器改周更换并记录数量趋势。\n"
     "③ 将 Orkin 服务报告纳入月度稽核输入，避免「主巡检没发现 = 没问题」\n"
     "  的盲区。", size=8.5)

# ================================================================= 30 FACILITY ATTRIBUTION
s = S(); title_bar(s, "设施（BD / 营建）扣分归因分析",
                   "综合 QA、门店自检、区经三类巡检", page())
fac_all = [it for it in ITEM if it["module"] == "设施"]
kpi(s, 0.55, 1.45, 2.9, 1.1, str(MA["设施"]["deduction"]), f"设施扣分（{NPRIM} 家口径）", "主巡检", CRIT, vsize=24)
kpi(s, 3.65, 1.45, 2.9, 1.1, f"10 / {PM['S']}", "主巡检设施关键项 / 占关键项", f"{10/PM['S']*100:.0f}%", CRIT, vsize=19)
kpi(s, 6.75, 1.45, 2.9, 1.1, "17 起", "全月 Sinks and Pipes S 项", "6 月 13 起", CRIT, vsize=24)
kpi(s, 9.85, 1.45, 2.9, 1.1, "11 家", "全月涉及门店", "6 月 10 家", CRIT, vsize=24)
band(s, 0.55, 2.75, 12.2, "问题类型与责任归因")
body(s, 0.55, 3.09, 12.2, 1.85,
     "问题类型：air gap 间距不足 / 不合规（最高频）、管道包覆层脱落或缺失、管道与滤芯间距过近、"
     "排水 backflow 与外溢、油脂阱管道异常、地漏连接不当。\n\n"
     f"设施模块主巡检扣分 {MA['设施']['deduction']} 分、覆盖 {MA['设施']['stores']}/{NPRIM} 家（{MA['设施']['coverage']}%），"
     f"且包含 10 个关键项（占全部关键项 {10/PM['S']*100:.0f}%）。\n"
     f"分数虽因申诉整体回升，但该类风险仍为 P0——本月 10 个 air gap 关键项 100% 申诉获批，扣分归零，"
     f"实物整改状态无系统记录。", size=9.5)
band(s, 0.55, 5.1, 12.2, "整改建议（P0 · BD / 营建）")
body(s, 0.55, 5.44, 12.2, 1.45,
     f"① BD 整改清单按门店逐项派单，明确 owner、预计完成日、复查人，8 月底前完工。\n"
     f"② 对 {len(REPEAT['repeat'])} 家连续两月复现门店（{'、'.join(sn(c) for c in REPEAT['repeat'][:6])} 等）优先排期。\n"
     f"③ 申诉获批与整改工单双向绑定：以「整改后复核」为由获批的，须留存完工验收单与整改前后照片。", size=9.5)

# ================================================================= 31 EQUIPMENT
s = S(); title_bar(s, "设备（Equipment）扣分项分析", "综合 QA、门店自检、区经三类巡检", page())
eq = [it for it in PITEMS if it["module"] == "设备维护"]
eq_all = [it for it in ITEM if it["module"] == "设备维护"]
kpi(s, 0.55, 1.45, 2.9, 1.1, f"{len(eq)} 项", "主巡检设备扣分项", f"6 月 0 项", WARN, vsize=24)
kpi(s, 3.65, 1.45, 2.9, 1.1, f"{MA['设备维护']['deduction']} 分", "设备扣分", None, WARN, vsize=24)
kpi(s, 6.75, 1.45, 2.9, 1.1, f"{MA['设备维护']['stores']} 家", "主巡检涉及门店", f"覆盖率 {MA['设备维护']['coverage']}%", WARN, vsize=24)
kpi(s, 9.85, 1.45, 2.9, 1.1, "G 项", "严重度", f"全月 {len(eq_all)} 项", GOOD, vsize=24)
band(s, 0.55, 2.75, 7.6, f"主巡检设备扣分项明细（共 {len(eq)} 项）")
table(s, 0.55, 3.12, 7.6, 1.4, ["门店", "问题描述", "严重度", "扣分"],
      [[sn(it["store_code"]), (it["description"] or "—").replace("\n"," ")[:46], it["severity"], it["deduction"]]
       for it in eq] or [["—", "本月无设备维护扣分项", "—", 0]],
      widths=[1.8,5.2,0.9,0.8], align_center=[2,3], fsize=8.5)
band(s, 8.4, 2.75, 4.35, "分析与建议")
body(s, 8.4, 3.09, 4.35, 3.6,
     f"设备维护本月出现 {len(eq)} 项主巡检扣分（6 月为 0），均为 G 项、扣分 {MA['设备维护']['deduction']} 分。\n\n"
     f"⚠ 高度集中：{len(eq)} 项全部为洗碗机故障（漏水 2、停机 1），涉及 100 Maiden Ln、41st & Lexington、40th & 10th，"
     f"提示该机型或维保周期存在共性问题，建议按机型批次排查。\n\n"
     f"⚠ 但需与供应链侧联动看待：7/28 共 {len([r for r in PJ if 'defrost' in (r['problem_description'] or '').lower()])} 起"
     f"后厨冰箱异常除霜导致物料报废"
     f"（货值 ${sum(val(r) for r in PJ if 'defrost' in (r['problem_description'] or '').lower()):,.2f}，见 Case C），"
     f"该损失记录在 PQNC 而非稽核扣分中。\n\n"
     f"建议：① 冷藏 / 冷冻设备加装温度记录仪并纳入日检；② 建立设备故障与 PQNC 损失的关联台账，"
     f"避免设备问题分散在两套系统中而低估真实风险。", size=9.5)

# ================================================================= 32 DIVIDER 04
s = S()
pic(s, "divider04", 0, 0.01, 8.83, 7.49)
section_head(s, "04", "Customer\nComplaint\n客户投诉", 9.89, 2.69, 2.92, 2.4)

# ================================================================= 33 COMPLAINT
s = S(); title_bar(s, "04  Customer Complaint / 客户投诉",
                   "2026 年 7 月 · 7 起 ※ 6 月 PPT 口径下 1–6 月累计仅 3 起（约 0.5 起/月），"
                   "与本月导出口径是否一致待客服确认，故本页不做 YTD 同比", page())
COMPLAINTS = [
    ("Jul 8, 11:06 AM",  "41st & Lexington", "The coconut milk has most definitely gone bad. "
                                             "The drink tastes sour and spoiled", "变质酸败"),
    ("Jul 12, 1:52 PM",  "28th & 6th",       "Tastes sour like the coconut milk is bad", "变质酸败"),
    ("Jul 14, 1:02 PM",  "37th & Broadway",  "The milk had curdled at the top where the orange "
                                             "flavor was added.", "变质酸败"),
    ("Jul 17, 2:42 PM",  "100 Maiden Ln",    "Strange chemical taste", "异味"),
    ("Jul 17, 8:33 PM",  "37th & Broadway",  "It taste so bad, i felt like i was drinking "
                                             "expired milk", "变质酸败"),
    ("Jul 24, 8:00 AM",  "15th & 3rd",       "Hair in food", "异物"),
    ("Jul 27, 12:42 PM", "221 Grand",        "There was a piece of plastic in my food", "异物"),
]
c_stores = {c[1] for c in COMPLAINTS}
c_spoil  = [c for c in COMPLAINTS if c[3] == "变质酸败"]
c_object = [c for c in COMPLAINTS if c[3] == "异物"]
cluster_stores = {ENR[r["pqnc_no"]]["store"] for r in JOINT}
kpi(s, 0.55, 1.4, 2.9, 1.15, str(len(COMPLAINTS)), "7 月食安类客诉",
    f"涉及 {len(c_stores)} 家门店", CRIT)
kpi(s, 3.65, 1.4, 2.9, 1.15, str(len(c_spoil)), "变质 / 酸败类", "全部指向奶类", CRIT)
kpi(s, 6.75, 1.4, 2.9, 1.15, str(len(c_object)), "异物类", "毛发 1 · 塑料 1", CRIT)
kpi(s, 9.85, 1.4, 2.9, 1.15, "11 天", "客诉早于内部工单", "7/8 首诉 → 7/19 首张变质 PQNC", CRIT, vsize=22)
band(s, 0.55, 2.72, 12.2, "事件明细 / Complaint Detail")
table(s, 0.55, 3.09, 12.2, 2.0,
      ["日期时间", "门店", "客诉原文 / Customer comment", "分类"],
      [[d, st, f"“{tx}”", cat] for d, st, tx, cat in COMPLAINTS],
      widths=[1.9,1.9,7.0,1.4], align_center=[0,3], fsize=8.5)
band(s, 0.55, 5.25, 6.0, "分析 / Analysis")
body(s, 0.55, 5.59, 6.0, 1.36,
     f"{len(COMPLAINTS)} 起中 {len(c_spoil)} 起为奶类变质酸败、1 起异味，{len(c_object)} 起异物（毛发、塑料）。\n"
     f"⚠ 检出滞后：客诉 7/8 即报椰奶变质，内部首张变质 PQNC 为 7/19、\n"
     f"集群 7/22 才爆发——「客户比工单早约两周发现同一问题」。\n"
     f"37th & Broadway 7/14、7/17 两次客诉，同店 7/21 才开 PQNC（滞后 7 天）。\n"
     f"41st & Lexington（7/8 首诉）与 15th & 3rd 全月「无」对应 PQNC，客诉未回流质量工单。\n"
     f"客诉 {len(c_stores)} 家中仅 {len(c_stores & cluster_stores)} 家落在 Case B 集群内。", size=8.5)
band(s, 6.75, 5.25, 6.0, "纠正措施 / Corrective Action（P1）")
body(s, 6.75, 5.59, 6.0, 1.36,
     "① 客诉与 PQNC 打通：奶类变质客诉自动触发同店同 SKU 留样 + 建单。\n"
     "② 以客诉为「早期指标」：同 SKU 跨店 2 起变质客诉即启动批次追溯\n"
     "   ——按本月时间线可提前约两周锁定 Cream-O-Land 批次。\n"
     "③ 41st & Lexington、15th & 3rd 补开工单并追溯对应批次。\n"
     "④ 异物类单独排查：15th & 3rd 员工毛发管控（发网/工帽）；\n"
     "   221 Grand 可颂供应商包装与门店加热出餐环节。", size=8.5)

# ================================================================= 34 DIVIDER 05
s = S()
pic(s, "divider05", 0, 0.09, 5.72, 7.41)
section_head(s, "05", "EHS\n环境健康安全", 7.68, 3.05, 4.6, 1.5)

# ================================================================= 35 EHS INCIDENTS
s = S(); title_bar(s, "2026 EHS 安全事故月度分析 / Monthly Safety Incident Analysis",
                   "工伤事故 Workplace Injury · YTD（截至 2026 年 7 月 / through July）", page())
kpi(s, 0.55, 1.4, 3.9, 1.2, "0", "7 月工伤事故 / Injuries", "无工伤", GOOD)
kpi(s, 4.65, 1.4, 3.9, 1.2, "1", "2026 YTD 累计", "仅 5 月 1 起", INFO)
kpi(s, 8.75, 1.4, 4.0, 1.2, "连续 2 月", "零工伤", "6 月 0 / 7 月 0", GOOD)
tb(s, 0.55, 2.95, 6.5, 0.25, "月度事故趋势（起）", size=9, bold=True, color=INK2)
chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, 0.5, 3.2, 7.4, 3.0,
      ["Jan 1月","Feb 2月","Mar 3月","Apr 4月","May 5月","Jun 6月","Jul 7月"],
      [("事故起数 Incidents", [0, 0, 0, 0, 1, 0, 0])], [C1], legend=False, lsize=9)
band(s, 8.2, 3.2, 4.55, "说明")
body(s, 8.2, 3.54, 4.55, 2.66,
     "7 月无工伤事故（0 起），6 月同为 0，连续两月保持零工伤。\n\n"
     "2026 YTD 累计 1 起，发生于 5 月。\n\n"
     "本月门店数由 18 家增至 21 家、巡检 94 次为历史最高，"
     "在业务扩张的同时保持零工伤。\n\n"
     "建议在月报中固化「事故起数 / 损工日数 / 高危作业检查次数」"
     "三项指标，便于跨月比较。", size=9.5)

# ================================================================= 36 EHS PROGRESS (placeholder)
s = S(); title_bar(s, "EHS 体系目前进程 / Current EHS System Progress", None, page())
placeholder(s, 0.55, 1.4, 12.2, 1.6, "EHS 体系进度待 EHS 团队更新",
    "以下为 6 月 PPT 的进度状态，7 月进展需 EHS 团队确认。",
    "特别是「4 家新店疏散图」——本月门店数已达 21 家，疏散图覆盖范围需同步更新。")
band(s, 0.55, 3.25, 6.0, "已完成 / Completed（6 月状态）")
body(s, 0.55, 3.62, 6.0, 3.05,
     "· 门店设备安全操作 SOP\n  SOP for safe operation of store equipment\n\n"
     "· 安全手册 Safety Manual\n\n"
     "· 基础安全培训 Basic Safety Training\n\n"
     "· 16 家门店疏散图 Evacuation plans for 16 stores", size=10)
band(s, 6.75, 3.25, 6.0, "进行中 / In Progress（待更新）")
body(s, 6.75, 3.62, 6.0, 3.05,
     "· 新店疏散图 Evacuation plans for new stores\n"
     "  ⚠ 6 月记录为「4 家新店」；截至 7 月底门店总数已达 21 家，\n"
     "  需确认疏散图覆盖缺口（21 − 16 = 5 家待完成）\n\n"
     "· 整体安全手册升级 Upgrade of the overall Safety Manual\n\n"
     "· 7 月新增进展：待填", size=10)

# ================================================================= 37 THANK YOU
s = S()
pic(s, "thanks", 0, 0, 8.75, 7.43)
tb(s, 9.30, 2.75, 3.60, 1.5, "Thank You！\n谢谢！", size=32, bold=True, color="223263", spacing=1.2)
tb(s, 9.32, 4.55, 3.60, 1.4, "北美质量管理部\n2026 年 7 月\n\n编制：曾翔宇", size=12, color=INK2, spacing=1.35)

prs.save(OUT_PATH)
print(f"[write] {OUT_PATH}")
print(f"[write] slides={len(prs.slides.__iter__.__self__._sldIdLst)}  size={OUT_PATH.stat().st_size:,} bytes")
print(f"[data ] PQNC July={N_JUL} (Jun={N_JUN})  QA={QA_N}/{QA_STORES} avg={QA_AVG}/{QA_AVG_ORG}")
